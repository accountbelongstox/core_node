#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Universal document parser for the File Processor MCP Server.

Extracted from main.py. Parses PDF/XMind/Word/Excel/PowerPoint/image/text/JSON/XML.
OCR runtime globals (ocr_manager, OCR_AVAILABLE) are read via function-local
imports from server to avoid capturing stale references and circular imports.
"""

import json
import logging
import mimetypes
import os
from pathlib import Path
from typing import Dict, Any, Optional

from package_manager import PackageManager
from document_converter import DocumentConverter
from pycore.pyfoundations.third_party import get_third_package_pypdf

logger = logging.getLogger(__name__)

# pypdf is resolved once at import (no circular dependency on server).
pypdf_lib = get_third_package_pypdf()


class DocumentParser:
    """Universal document parser with read and write capabilities"""

    def __init__(self):
        self.supported_formats = {
            '.pdf': self.parse_pdf,
            '.xmind': self.parse_xmind,
            '.docx': self.parse_word,
            '.doc': self.parse_word,
            '.xlsx': self.parse_excel,
            '.xls': self.parse_excel,
            '.pptx': self.parse_powerpoint,
            '.ppt': self.parse_powerpoint,
            '.txt': self.parse_text,
            '.md': self.parse_text,
            '.json': self.parse_json,
            '.xml': self.parse_xml,
            '.jpg': self.parse_image,
            '.jpeg': self.parse_image,
            '.png': self.parse_image,
            '.bmp': self.parse_image,
            '.tiff': self.parse_image,
        }

        # Initialize converter
        self.converter = DocumentConverter()

    def detect_file_type(self, file_path: str) -> Optional[str]:
        """Detect file type from extension and MIME type"""
        file_path = Path(file_path)
        extension = file_path.suffix.lower()

        if extension in self.supported_formats:
            return extension

        # Try MIME type detection
        mime_type, _ = mimetypes.guess_type(str(file_path))
        if mime_type:
            mime_mapping = {
                'application/pdf': '.pdf',
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
                'text/plain': '.txt',
                'text/markdown': '.md',
                'application/json': '.json',
                'application/xml': '.xml',
                'text/xml': '.xml',
                'image/jpeg': '.jpg',
                'image/png': '.png',
                'image/bmp': '.bmp',
                'image/tiff': '.tiff',
            }
            return mime_mapping.get(mime_type)

        return None

    def parse_pdf(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse PDF file"""
        try:
            # Ensure PDF packages are installed
            if not PackageManager.ensure_packages('pdf'):
                raise ImportError("PDF parsing packages not available")

            if pypdf_lib is None:
                raise ImportError("pypdf not available")
            import pdfplumber

            result = {
                "type": "pdf",
                "file_path": file_path,
                "metadata": {},
                "content": {
                    "text": "",
                    "pages": [],
                    "tables": [],
                    "images": []
                }
            }

            # Extract text with pypdf
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf_lib.PdfReader(file)
                result["metadata"] = {
                    "num_pages": len(pdf_reader.pages),
                    "title": pdf_reader.metadata.get('/Title', '') if pdf_reader.metadata else '',
                    "author": pdf_reader.metadata.get('/Author', '') if pdf_reader.metadata else '',
                    "creator": pdf_reader.metadata.get('/Creator', '') if pdf_reader.metadata else '',
                }

                full_text = []
                for i, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text() or ""
                    full_text.append(page_text)
                    result["content"]["pages"].append({
                        "page": i + 1,
                        "text": page_text
                    })

                result["content"]["text"] = "\n\n".join(full_text)

            # Extract tables with pdfplumber
            try:
                with pdfplumber.open(file_path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        tables = page.extract_tables()
                        for j, table in enumerate(tables):
                            if table:
                                result["content"]["tables"].append({
                                    "page": i + 1,
                                    "table_index": j,
                                    "data": table
                                })
            except Exception as e:
                logger.warning(f"Table extraction failed: {e}")

            return result

        except Exception as e:
            logger.error(f"PDF parsing failed: {e}")
            return {"error": str(e), "type": "pdf", "file_path": file_path}

    def parse_xmind(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse XMind file"""
        try:
            # Ensure XMind packages are installed
            if not PackageManager.ensure_packages('xmind'):
                raise ImportError("XMind parsing packages not available")

            from xmindparser import xmind_to_dict

            result = {
                "type": "xmind",
                "file_path": file_path,
                "metadata": {
                    "file_size": os.path.getsize(file_path)
                },
                "content": {}
            }

            # Parse XMind file
            xmind_data = xmind_to_dict(file_path)
            result["content"] = xmind_data

            # Extract plain text for AI
            def extract_text_from_xmind(data, level=0):
                text_parts = []
                indent = "  " * level

                if isinstance(data, list):
                    for item in data:
                        text_parts.extend(extract_text_from_xmind(item, level))
                elif isinstance(data, dict):
                    # Handle topic structure
                    if 'topic' in data:
                        topic_data = data['topic']
                        if 'title' in topic_data:
                            text_parts.append(f"{indent}- {topic_data['title']}")
                        if 'topics' in topic_data:
                            text_parts.extend(extract_text_from_xmind(topic_data['topics'], level + 1))
                    else:
                        for key, value in data.items():
                            if key == 'title' and isinstance(value, str):
                                text_parts.append(f"{indent}- {value}")
                            elif isinstance(value, (list, dict)):
                                text_parts.extend(extract_text_from_xmind(value, level + 1))

                return text_parts

            text_content = extract_text_from_xmind(xmind_data)
            result["content"]["extracted_text"] = "\n".join(text_content)

            return result

        except Exception as e:
            logger.error(f"XMind parsing failed: {e}")
            return {"error": str(e), "type": "xmind", "file_path": file_path}

    def parse_word(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse Word document"""
        try:
            # Ensure Office packages are installed
            if not PackageManager.ensure_packages('office'):
                raise ImportError("Office parsing packages not available")

            from docx import Document

            result = {
                "type": "word",
                "file_path": file_path,
                "metadata": {},
                "content": {
                    "text": "",
                    "paragraphs": [],
                    "tables": [],
                    "styles": []
                }
            }

            # Parse document
            doc = Document(file_path)

            # Extract metadata
            core_props = doc.core_properties
            result["metadata"] = {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
                "subject": core_props.subject or "",
                "keywords": core_props.keywords or ""
            }

            # Extract content
            full_text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    result["content"]["paragraphs"].append({
                        "text": paragraph.text,
                        "style": paragraph.style.name if paragraph.style else "Normal"
                    })
                    full_text.append(paragraph.text)

            # Extract tables
            for i, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)

                result["content"]["tables"].append({
                    "table_index": i,
                    "data": table_data
                })

            result["content"]["text"] = "\n\n".join(full_text)

            return result

        except Exception as e:
            logger.error(f"Word parsing failed: {e}")
            return {"error": str(e), "type": "word", "file_path": file_path}

    def parse_excel(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse Excel spreadsheet"""
        try:
            # Ensure Office packages are installed
            if not PackageManager.ensure_packages('office'):
                raise ImportError("Office parsing packages not available")

            import openpyxl

            result = {
                "type": "excel",
                "file_path": file_path,
                "metadata": {},
                "content": {
                    "worksheets": [],
                    "summary": ""
                }
            }

            # Load workbook
            workbook = openpyxl.load_workbook(file_path, read_only=True)

            # Extract metadata
            result["metadata"] = {
                "worksheet_count": len(workbook.worksheets),
                "worksheet_names": [ws.title for ws in workbook.worksheets]
            }

            # Extract content from each worksheet
            summary_parts = []
            for worksheet in workbook.worksheets:
                ws_data = {
                    "name": worksheet.title,
                    "max_row": worksheet.max_row,
                    "max_column": worksheet.max_column,
                    "data": []
                }

                # Read data (limit to reasonable size)
                max_rows = min(worksheet.max_row, 1000)
                max_cols = min(worksheet.max_column, 50)

                for row in worksheet.iter_rows(min_row=1, max_row=max_rows,
                                             min_col=1, max_col=max_cols, values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(cell.strip() for cell in row_data):  # Skip empty rows
                        ws_data["data"].append(row_data)

                result["content"]["worksheets"].append(ws_data)
                summary_parts.append(f"Worksheet '{worksheet.title}': {len(ws_data['data'])} rows of data")

            result["content"]["summary"] = "\n".join(summary_parts)

            return result

        except Exception as e:
            logger.error(f"Excel parsing failed: {e}")
            return {"error": str(e), "type": "excel", "file_path": file_path}

    def parse_powerpoint(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse PowerPoint presentation"""
        try:
            # Ensure Office packages are installed
            if not PackageManager.ensure_packages('office'):
                raise ImportError("Office parsing packages not available")

            from pptx import Presentation

            result = {
                "type": "powerpoint",
                "file_path": file_path,
                "metadata": {},
                "content": {
                    "text": "",
                    "slides": [],
                    "notes": []
                }
            }

            # Load presentation
            prs = Presentation(file_path)

            # Extract metadata
            result["metadata"] = {
                "slide_count": len(prs.slides),
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height
            }

            # Extract content
            full_text = []
            for i, slide in enumerate(prs.slides):
                slide_data = {
                    "slide_number": i + 1,
                    "title": "",
                    "content": [],
                    "notes": ""
                }

                # Extract slide content
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if shape.placeholder_format and shape.placeholder_format.type == 1:  # Title
                            slide_data["title"] = shape.text
                        else:
                            slide_data["content"].append(shape.text)
                        full_text.append(shape.text)

                # Extract notes
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    slide_data["notes"] = notes_text
                    if notes_text.strip():
                        full_text.append(f"Notes: {notes_text}")

                result["content"]["slides"].append(slide_data)

            result["content"]["text"] = "\n\n".join(full_text)

            return result

        except Exception as e:
            logger.error(f"PowerPoint parsing failed: {e}")
            return {"error": str(e), "type": "powerpoint", "file_path": file_path}

    def parse_image(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse image using advanced OCR engines"""
        try:
            result = {
                "type": "image",
                "file_path": file_path,
                "metadata": {},
                "content": {
                    "text": "",
                    "confidence": 0,
                    "words": [],
                    "lines": [],
                    "ocr_results": []
                }
            }

            # Load image metadata
            try:
                from PIL import Image
                with Image.open(file_path) as img:
                    result["metadata"] = {
                        "width": img.width,
                        "height": img.height,
                        "mode": img.mode,
                        "format": img.format
                    }
            except Exception as e:
                logger.warning(f"Could not load image metadata: {e}")

            # OCR runtime globals are read function-locally to avoid stale captures.
            from server import OCR_AVAILABLE, ocr_manager

            # Use advanced OCR engines if available
            if OCR_AVAILABLE and ocr_manager:
                # Get OCR engine preference
                ocr_engine = kwargs.get('ocr_engine', 'auto')
                use_fallback = kwargs.get('use_fallback', True)

                if ocr_engine == 'auto' or use_fallback:
                    # Use fallback mechanism
                    engines = kwargs.get('engines', None)
                    ocr_result = ocr_manager.recognize_with_fallback(file_path, engines)
                else:
                    # Use specific engine
                    ocr_result = ocr_manager.recognize(file_path, ocr_engine)

                if ocr_result.success:
                    result["content"]["text"] = ocr_result.text
                    result["content"]["confidence"] = ocr_result.confidence
                    result["content"]["words"] = ocr_result.words
                    result["content"]["ocr_results"].append(ocr_result.to_dict())
                    logger.info(f"OCR successful with {ocr_result.provider}")
                else:
                    logger.warning(f"Advanced OCR failed: {ocr_result.error}")
                    # Fallback to pytesseract
                    return self._fallback_ocr(file_path, result, **kwargs)
            else:
                # Fallback to pytesseract
                return self._fallback_ocr(file_path, result, **kwargs)

            return result

        except Exception as e:
            logger.error(f"Image OCR failed: {e}")
            return {"error": str(e), "type": "image", "file_path": file_path}

    def _fallback_ocr(self, file_path: str, result: Dict[str, Any], **kwargs) -> Dict[str, Any]:
        """Fallback OCR using pytesseract"""
        try:
            # Ensure OCR packages are installed
            if not PackageManager.ensure_packages('ocr'):
                raise ImportError("OCR parsing packages not available")

            import pytesseract
            from PIL import Image

            logger.info("Using pytesseract fallback OCR")

            # Load and analyze image
            with Image.open(file_path) as img:
                # OCR extraction
                ocr_confidence = kwargs.get('ocr_confidence', 30)
                ocr_data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)

                # Filter confident text
                confident_text = []
                confidences = []
                words = []

                for i, confidence in enumerate(ocr_data['conf']):
                    if int(confidence) > ocr_confidence:
                        text = ocr_data['text'][i].strip()
                        if text:
                            confident_text.append(text)
                            confidences.append(int(confidence))
                            words.append({
                                "text": text,
                                "confidence": int(confidence),
                                "bbox": {
                                    "left": ocr_data['left'][i],
                                    "top": ocr_data['top'][i],
                                    "width": ocr_data['width'][i],
                                    "height": ocr_data['height'][i]
                                }
                            })

                result["content"]["text"] = " ".join(confident_text)
                result["content"]["confidence"] = sum(confidences) / len(confidences) if confidences else 0
                result["content"]["words"] = words
                result["content"]["ocr_results"].append({
                    "provider": "pytesseract",
                    "success": True,
                    "text": result["content"]["text"],
                    "confidence": result["content"]["confidence"]
                })

            return result

        except Exception as e:
            logger.error(f"Fallback OCR failed: {e}")
            result["content"]["ocr_results"].append({
                "provider": "pytesseract",
                "success": False,
                "error": str(e)
            })
            return result

    def parse_text(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse text file"""
        try:
            result = {
                "type": "text",
                "file_path": file_path,
                "metadata": {},
                "content": {
                    "text": "",
                    "lines": 0,
                    "encoding": "utf-8"
                }
            }

            # Try different encodings
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
            content = None

            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    result["content"]["encoding"] = encoding
                    break
                except UnicodeDecodeError:
                    continue

            if content is None:
                # Fallback to binary read
                with open(file_path, 'rb') as f:
                    content = f.read().decode('utf-8', errors='ignore')
                result["content"]["encoding"] = "utf-8 (with errors ignored)"

            result["content"]["text"] = content
            result["content"]["lines"] = len(content.splitlines())
            result["metadata"]["file_size"] = os.path.getsize(file_path)

            return result

        except Exception as e:
            logger.error(f"Text parsing failed: {e}")
            return {"error": str(e), "type": "text", "file_path": file_path}

    def parse_json(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse JSON file"""
        try:
            result = {
                "type": "json",
                "file_path": file_path,
                "metadata": {},
                "content": {}
            }

            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)

            result["content"] = data
            result["metadata"]["file_size"] = os.path.getsize(file_path)

            return result

        except Exception as e:
            logger.error(f"JSON parsing failed: {e}")
            return {"error": str(e), "type": "json", "file_path": file_path}

    def parse_xml(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse XML file"""
        try:
            # Ensure XML packages are installed
            if not PackageManager.ensure_packages('xml'):
                raise ImportError("XML parsing packages not available")

            import xml.etree.ElementTree as ET

            result = {
                "type": "xml",
                "file_path": file_path,
                "metadata": {},
                "content": {
                    "text": "",
                    "structure": {}
                }
            }

            # Parse XML
            tree = ET.parse(file_path)
            root = tree.getroot()

            def xml_to_dict(element):
                result = {}
                if element.text and element.text.strip():
                    result['text'] = element.text.strip()
                if element.attrib:
                    result['attributes'] = element.attrib

                children = {}
                for child in element:
                    child_data = xml_to_dict(child)
                    if child.tag in children:
                        if not isinstance(children[child.tag], list):
                            children[child.tag] = [children[child.tag]]
                        children[child.tag].append(child_data)
                    else:
                        children[child.tag] = child_data

                if children:
                    result['children'] = children

                return result

            result["content"]["structure"] = {root.tag: xml_to_dict(root)}

            # Extract all text content
            def extract_text(element):
                text_parts = []
                if element.text and element.text.strip():
                    text_parts.append(element.text.strip())
                for child in element:
                    text_parts.extend(extract_text(child))
                return text_parts

            all_text = extract_text(root)
            result["content"]["text"] = "\n".join(all_text)
            result["metadata"]["file_size"] = os.path.getsize(file_path)

            return result

        except Exception as e:
            logger.error(f"XML parsing failed: {e}")
            return {"error": str(e), "type": "xml", "file_path": file_path}

    def parse_file(self, file_path: str, output_format: str = "json", **kwargs) -> Dict[str, Any]:
        """Parse any supported file format"""
        try:
            # Validate file exists
            if not os.path.exists(file_path):
                return {"error": f"File not found: {file_path}"}

            # Detect file type
            file_type = self.detect_file_type(file_path)
            if not file_type:
                return {"error": f"Unsupported file format: {file_path}"}

            # Parse file
            parser_func = self.supported_formats[file_type]
            result = parser_func(file_path, **kwargs)

            # Add general metadata
            if "metadata" not in result:
                result["metadata"] = {}

            result["metadata"].update({
                "file_name": os.path.basename(file_path),
                "file_size": os.path.getsize(file_path),
                "detected_type": file_type,
                "parser_version": "1.0"
            })

            return result

        except Exception as e:
            logger.error(f"File parsing failed: {e}")
            return {"error": str(e), "file_path": file_path}

# TODO: A redundant sibling DocumentParserWithTextPositions exists under
# pycore/pyutils/mcp/file_processing/... and is not consolidated here because it is
# too cross-cutting (different package/layout). Consolidate in a future pass.
