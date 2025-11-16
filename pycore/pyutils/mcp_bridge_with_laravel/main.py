#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
File Processor MCP Server (文件解析器)
Universal file and document processor for AI-readable content extraction

Supported formats:
- PDF files (text and OCR)
- XMind files (mind maps)
- Word documents (.docx, .doc)
- Excel spreadsheets (.xlsx, .xls)
- PowerPoint presentations (.pptx, .ppt)
- Images (OCR extraction and pixel analysis)
- Text files

Features:
- Auto-detect and install required packages
- Smart content extraction with structure preservation
- Multiple output formats (JSON, XML, Markdown)
- Image pixel analysis with color matrix extraction
- Comprehensive error handling and debugging
- AI-optimized content formatting
"""

import asyncio
import json
import logging
import os
import sys
import subprocess
import importlib
import time
from pathlib import Path
from typing import Dict, Any, Optional, List, Union, Tuple
import mimetypes
import threading

# FastMCP framework
from pycore.pyfoundations.third_party import get_third_package_FastMCP

FastMCP = get_third_package_FastMCP()

# Import OCR engines and queue system
try:
    from ocr_engines import ocr_manager, OCRResult
    from ocr_queue_system import ocr_queue, TaskPriority
    from ocr_config import OCRLimits, ProcessingConfig
    from image_processor import SmartImageProcessor
    from pdf_processor import PDFProcessor
    OCR_AVAILABLE = True
except ImportError as e:
    logger.warning(f"OCR engines not available: {e}")
    OCR_AVAILABLE = False
    ocr_manager = None
    ocr_queue = None

# Import code scanner
try:
    from code_scanner import scanner as code_scanner
    CODE_SCANNER_AVAILABLE = True
except ImportError as e:
    logger.warning(f"Code scanner not available: {e}")
    CODE_SCANNER_AVAILABLE = False
    code_scanner = None

# Initialize MCP server
mcp = FastMCP("FileProcessor")

# Configure logging
# IMPORTANT: Use stderr for logging to avoid interfering with MCP JSON-RPC on stdout
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] %(message)s',
    handlers=[
        logging.StreamHandler(sys.stderr)  # Use stderr instead of stdout for MCP compatibility
    ]
)
logger = logging.getLogger(__name__)

class PackageManager:
    """Smart package installation manager with non-blocking initialization"""

    PACKAGE_MAPPING = {
        'pdf': ['PyPDF2', 'pdfplumber', 'pytesseract'],
        'xmind': ['xmindparser'],
        'office': ['python-docx', 'openpyxl', 'python-pptx'],
        'ocr': ['pytesseract', 'Pillow', 'requests'],
        'xml': ['dicttoxml', 'lxml'],
        'markdown': ['markdown', 'html2text'],
        'conversion': ['pandoc', 'weasyprint', 'pdfkit'],
        'html': ['beautifulsoup4', 'html2text', 'markdownify'],
        'paddle_ocr': ['paddlepaddle', 'paddleocr'],
        'cnocr': ['cnocr[ort-cpu]']
    }

    _installation_status = {}
    _installation_lock = threading.Lock()

    @staticmethod
    def install_package(package_name: str, timeout: int = 60) -> bool:
        """Install a Python package with timeout"""
        with PackageManager._installation_lock:
            # Check cache first
            if package_name in PackageManager._installation_status:
                return PackageManager._installation_status[package_name]

            try:
                importlib.import_module(package_name.replace('-', '_').replace('python_', ''))
                logger.info(f"[OK] {package_name} is already installed")
                PackageManager._installation_status[package_name] = True
                return True
            except ImportError:
                logger.info(f"[INFO] Installing {package_name}...")
                try:
                    process = subprocess.Popen([
                        sys.executable, "-m", "pip", "install", package_name
                    ], stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)

                    try:
                        stdout, stderr = process.communicate(timeout=timeout)
                        if process.returncode == 0:
                            logger.info(f"[OK] {package_name} installed successfully")
                            PackageManager._installation_status[package_name] = True
                            return True
                        else:
                            logger.error(f"[ERROR] Failed to install {package_name}: {stderr}")
                            PackageManager._installation_status[package_name] = False
                            return False
                    except subprocess.TimeoutExpired:
                        process.kill()
                        logger.error(f"[ERROR] {package_name} installation timed out")
                        PackageManager._installation_status[package_name] = False
                        return False

                except Exception as e:
                    logger.error(f"[ERROR] Failed to install {package_name}: {e}")
                    PackageManager._installation_status[package_name] = False
                    return False

    @staticmethod
    def ensure_packages(format_type: str, async_install: bool = False) -> bool:
        """Ensure required packages for format type are installed"""
        packages = PackageManager.PACKAGE_MAPPING.get(format_type, [])

        if async_install:
            # Start installation in background
            def install_async():
                for package in packages:
                    PackageManager.install_package(package)

            thread = threading.Thread(target=install_async, daemon=True)
            thread.start()
            return True
        else:
            success_count = 0
            for package in packages:
                if PackageManager.install_package(package):
                    success_count += 1
                else:
                    logger.warning(f"[WARN] Optional package {package} not installed")

            # Return True if at least half of the packages are installed
            return success_count >= len(packages) // 2 + 1 if packages else True

    @staticmethod
    def initialize_all_packages():
        """Initialize all packages in background (non-blocking)"""
        def init_packages():
            logger.info("[INFO] Starting background package initialization...")
            for format_type in PackageManager.PACKAGE_MAPPING:
                PackageManager.ensure_packages(format_type, async_install=False)
            logger.info("[INFO] Package initialization completed")

        thread = threading.Thread(target=init_packages, daemon=True)
        thread.start()
        logger.info("[INFO] Package initialization started in background")

class DocumentConverter:
    """Document format converter and writer"""

    def __init__(self):
        self.conversion_mappings = {
            'html_to_markdown': self.html_to_markdown,
            'html_to_pdf': self.html_to_pdf,
            'html_to_docx': self.html_to_docx,
            'markdown_to_html': self.markdown_to_html,
            'markdown_to_pdf': self.markdown_to_pdf,
            'text_to_html': self.text_to_html,
            'json_to_excel': self.json_to_excel,
            'excel_to_json': self.excel_to_json
        }

    def html_to_markdown(self, html_content: str, **kwargs) -> str:
        """Convert HTML to Markdown"""
        try:
            if not PackageManager.ensure_packages('html'):
                raise ImportError("HTML conversion packages not available")

            import html2text
            h = html2text.HTML2Text()
            h.ignore_links = kwargs.get('ignore_links', False)
            h.ignore_images = kwargs.get('ignore_images', False)
            h.body_width = kwargs.get('body_width', 0)
            return h.handle(html_content)
        except Exception as e:
            logger.error(f"HTML to Markdown conversion failed: {e}")
            raise

    def html_to_pdf(self, html_content: str, output_path: str, **kwargs) -> bool:
        """Convert HTML to PDF"""
        try:
            if not PackageManager.ensure_packages('conversion'):
                raise ImportError("PDF conversion packages not available")

            try:
                import pdfkit
                options = {
                    'page-size': kwargs.get('page_size', 'A4'),
                    'encoding': kwargs.get('encoding', 'UTF-8'),
                    'margin-top': kwargs.get('margin_top', '0.75in'),
                    'margin-right': kwargs.get('margin_right', '0.75in'),
                    'margin-bottom': kwargs.get('margin_bottom', '0.75in'),
                    'margin-left': kwargs.get('margin_left', '0.75in'),
                }
                pdfkit.from_string(html_content, output_path, options=options)
                return True
            except ImportError:
                # Fallback to weasyprint
                import weasyprint
                html_doc = weasyprint.HTML(string=html_content)
                html_doc.write_pdf(output_path)
                return True

        except Exception as e:
            logger.error(f"HTML to PDF conversion failed: {e}")
            return False

    def html_to_docx(self, html_content: str, output_path: str, **kwargs) -> bool:
        """Convert HTML to Word document"""
        try:
            if not PackageManager.ensure_packages('html'):
                raise ImportError("HTML conversion packages not available")

            from bs4 import BeautifulSoup

            # First convert HTML to markdown, then to docx
            markdown_content = self.html_to_markdown(html_content)
            return self.markdown_to_docx(markdown_content, output_path, **kwargs)

        except Exception as e:
            logger.error(f"HTML to DOCX conversion failed: {e}")
            return False

    def markdown_to_html(self, markdown_content: str, **kwargs) -> str:
        """Convert Markdown to HTML"""
        try:
            if not PackageManager.ensure_packages('markdown'):
                raise ImportError("Markdown conversion packages not available")

            import markdown
            md = markdown.Markdown(extensions=kwargs.get('extensions', ['tables', 'fenced_code']))
            return md.convert(markdown_content)

        except Exception as e:
            logger.error(f"Markdown to HTML conversion failed: {e}")
            raise

    def markdown_to_pdf(self, markdown_content: str, output_path: str, **kwargs) -> bool:
        """Convert Markdown to PDF"""
        try:
            # Convert markdown to HTML first, then to PDF
            html_content = self.markdown_to_html(markdown_content, **kwargs)
            return self.html_to_pdf(html_content, output_path, **kwargs)

        except Exception as e:
            logger.error(f"Markdown to PDF conversion failed: {e}")
            return False

    def markdown_to_docx(self, markdown_content: str, output_path: str, **kwargs) -> bool:
        """Convert Markdown to Word document"""
        try:
            if not PackageManager.ensure_packages('office'):
                raise ImportError("Office conversion packages not available")

            from docx import Document
            from docx.shared import Inches

            doc = Document()
            lines = markdown_content.split('\n')

            for line in lines:
                line = line.strip()
                if not line:
                    continue

                if line.startswith('# '):
                    doc.add_heading(line[2:], level=1)
                elif line.startswith('## '):
                    doc.add_heading(line[3:], level=2)
                elif line.startswith('### '):
                    doc.add_heading(line[4:], level=3)
                elif line.startswith('- ') or line.startswith('* '):
                    doc.add_paragraph(line[2:], style='List Bullet')
                else:
                    doc.add_paragraph(line)

            doc.save(output_path)
            return True

        except Exception as e:
            logger.error(f"Markdown to DOCX conversion failed: {e}")
            return False

    def text_to_html(self, text_content: str, **kwargs) -> str:
        """Convert plain text to HTML"""
        try:
            lines = text_content.split('\n')
            html_lines = []

            for line in lines:
                if not line.strip():
                    html_lines.append('<br>')
                else:
                    html_lines.append(f'<p>{line}</p>')

            title = kwargs.get('title', 'Document')

            html_template = f"""
            <!DOCTYPE html>
            <html lang="en">
            <head>
                <meta charset="UTF-8">
                <meta name="viewport" content="width=device-width, initial-scale=1.0">
                <title>{title}</title>
                <style>
                    body {{ font-family: Arial, sans-serif; line-height: 1.6; margin: 40px; }}
                    p {{ margin-bottom: 1em; }}
                </style>
            </head>
            <body>
                {''.join(html_lines)}
            </body>
            </html>
            """

            return html_template.strip()

        except Exception as e:
            logger.error(f"Text to HTML conversion failed: {e}")
            raise

    def json_to_excel(self, json_data: Dict[str, Any], output_path: str, **kwargs) -> bool:
        """Convert JSON data to Excel file"""
        try:
            if not PackageManager.ensure_packages('office'):
                raise ImportError("Office conversion packages not available")

            import openpyxl
            from openpyxl import Workbook

            wb = Workbook()
            ws = wb.active
            ws.title = kwargs.get('sheet_name', 'Data')

            def flatten_json(data, parent_key='', sep='.'):
                items = []
                if isinstance(data, dict):
                    for k, v in data.items():
                        new_key = f"{parent_key}{sep}{k}" if parent_key else k
                        if isinstance(v, (dict, list)):
                            items.extend(flatten_json(v, new_key, sep).items())
                        else:
                            items.append((new_key, v))
                elif isinstance(data, list):
                    for i, v in enumerate(data):
                        new_key = f"{parent_key}[{i}]"
                        if isinstance(v, (dict, list)):
                            items.extend(flatten_json(v, new_key, sep).items())
                        else:
                            items.append((new_key, v))
                return dict(items)

            flat_data = flatten_json(json_data)

            # Write headers
            headers = list(flat_data.keys())
            for col, header in enumerate(headers, 1):
                ws.cell(row=1, column=col, value=header)

            # Write data
            values = list(flat_data.values())
            for col, value in enumerate(values, 1):
                ws.cell(row=2, column=col, value=str(value))

            wb.save(output_path)
            return True

        except Exception as e:
            logger.error(f"JSON to Excel conversion failed: {e}")
            return False

    def excel_to_json(self, excel_path: str, **kwargs) -> Dict[str, Any]:
        """Convert Excel file to JSON data"""
        try:
            if not PackageManager.ensure_packages('office'):
                raise ImportError("Office conversion packages not available")

            import openpyxl

            wb = openpyxl.load_workbook(excel_path)
            result = {}

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_data = []

                # Get headers from first row
                headers = [cell.value for cell in ws[1]]

                # Get data rows
                for row in ws.iter_rows(min_row=2, values_only=True):
                    if any(cell is not None for cell in row):
                        row_data = {}
                        for header, cell in zip(headers, row):
                            if header:
                                row_data[header] = cell
                        if row_data:
                            sheet_data.append(row_data)

                result[sheet_name] = sheet_data

            return result

        except Exception as e:
            logger.error(f"Excel to JSON conversion failed: {e}")
            raise

    def convert_format(
        self,
        input_data: Union[str, Dict[str, Any]],
        from_format: str,
        to_format: str,
        output_path: Optional[str] = None,
        **kwargs
    ) -> Union[str, Dict[str, Any], bool]:
        """Universal format converter"""
        try:
            conversion_key = f"{from_format}_to_{to_format}"

            if conversion_key not in self.conversion_mappings:
                raise ValueError(f"Conversion from {from_format} to {to_format} not supported")

            converter_func = self.conversion_mappings[conversion_key]

            if output_path:
                # File output conversion
                if from_format == 'html' and to_format in ['pdf', 'docx']:
                    return converter_func(input_data, output_path, **kwargs)
                elif from_format == 'markdown' and to_format in ['pdf', 'docx']:
                    return converter_func(input_data, output_path, **kwargs)
                elif from_format == 'json' and to_format == 'excel':
                    return converter_func(input_data, output_path, **kwargs)
                else:
                    # String conversion, then save to file
                    result = converter_func(input_data, **kwargs)
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(result)
                    return True
            else:
                # String/data conversion
                return converter_func(input_data, **kwargs)

        except Exception as e:
            logger.error(f"Format conversion failed: {e}")
            raise

class DocumentParser:
    """Universal document parser with read and write capabilities"""

    def __init__(self):
        self.supported_formats = {
            '.pdf': self.parse_pdf,
            '.xmind': self.parse_xmind,
            '.docx': self.parse_word,
            '.doc': self.parse_word,
            '.xlsx': self.parse_excel,
            '.xls': self.parse_excel,
            '.pptx': self.parse_powerpoint,
            '.ppt': self.parse_powerpoint,
            '.txt': self.parse_text,
            '.md': self.parse_text,
            '.json': self.parse_json,
            '.xml': self.parse_xml,
            '.jpg': self.parse_image,
            '.jpeg': self.parse_image,
            '.png': self.parse_image,
            '.bmp': self.parse_image,
            '.tiff': self.parse_image,
        }

        # Initialize converter
        self.converter = DocumentConverter()

    def detect_file_type(self, file_path: str) -> Optional[str]:
        """Detect file type from extension and MIME type"""
        file_path = Path(file_path)
        extension = file_path.suffix.lower()

        if extension in self.supported_formats:
            return extension

        # Try MIME type detection
        mime_type, _ = mimetypes.guess_type(str(file_path))
        if mime_type:
            mime_mapping = {
                'application/pdf': '.pdf',
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
                'text/plain': '.txt',
                'text/markdown': '.md',
                'application/json': '.json',
                'application/xml': '.xml',
                'text/xml': '.xml',
                'image/jpeg': '.jpg',
                'image/png': '.png',
                'image/bmp': '.bmp',
                'image/tiff': '.tiff',
            }
            return mime_mapping.get(mime_type)

        return None

    def parse_pdf(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse PDF file"""
        try:
            # Ensure PDF packages are installed
            if not PackageManager.ensure_packages('pdf'):
                raise ImportError("PDF parsing packages not available")

            import PyPDF2
            import pdfplumber

            result = {
                "type": "pdf",
                "file_path": file_path,
                "metadata": {},
                "content": {
                    "text": "",
                    "pages": [],
                    "tables": [],
                    "images": []
                }
            }

            # Extract text with PyPDF2
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                result["metadata"] = {
                    "num_pages": len(pdf_reader.pages),
                    "title": pdf_reader.metadata.get('/Title', '') if pdf_reader.metadata else '',
                    "author": pdf_reader.metadata.get('/Author', '') if pdf_reader.metadata else '',
                    "creator": pdf_reader.metadata.get('/Creator', '') if pdf_reader.metadata else '',
                }

                full_text = []
                for i, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text() or ""
                    full_text.append(page_text)
                    result["content"]["pages"].append({
                        "page": i + 1,
                        "text": page_text
                    })

                result["content"]["text"] = "\n\n".join(full_text)

            # Extract tables with pdfplumber
            try:
                with pdfplumber.open(file_path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        tables = page.extract_tables()
                        for j, table in enumerate(tables):
                            if table:
                                result["content"]["tables"].append({
                                    "page": i + 1,
                                    "table_index": j,
                                    "data": table
                                })
            except Exception as e:
                logger.warning(f"Table extraction failed: {e}")

            return result

        except Exception as e:
            logger.error(f"PDF parsing failed: {e}")
            return {"error": str(e), "type": "pdf", "file_path": file_path}

    def parse_xmind(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse XMind file"""
        try:
            # Ensure XMind packages are installed
            if not PackageManager.ensure_packages('xmind'):
                raise ImportError("XMind parsing packages not available")

            from xmindparser import xmind_to_dict

            result = {
                "type": "xmind",
                "file_path": file_path,
                "metadata": {
                    "file_size": os.path.getsize(file_path)
                },
                "content": {}
            }

            # Parse XMind file
            xmind_data = xmind_to_dict(file_path)
            result["content"] = xmind_data

            # Extract plain text for AI
            def extract_text_from_xmind(data, level=0):
                text_parts = []
                indent = "  " * level

                if isinstance(data, list):
                    for item in data:
                        text_parts.extend(extract_text_from_xmind(item, level))
                elif isinstance(data, dict):
                    # Handle topic structure
                    if 'topic' in data:
                        topic_data = data['topic']
                        if 'title' in topic_data:
                            text_parts.append(f"{indent}- {topic_data['title']}")
                        if 'topics' in topic_data:
                            text_parts.extend(extract_text_from_xmind(topic_data['topics'], level + 1))
                    else:
                        for key, value in data.items():
                            if key == 'title' and isinstance(value, str):
                                text_parts.append(f"{indent}- {value}")
                            elif isinstance(value, (list, dict)):
                                text_parts.extend(extract_text_from_xmind(value, level + 1))

                return text_parts

            text_content = extract_text_from_xmind(xmind_data)
            result["content"]["extracted_text"] = "\n".join(text_content)

            return result

        except Exception as e:
            logger.error(f"XMind parsing failed: {e}")
            return {"error": str(e), "type": "xmind", "file_path": file_path}

    def parse_word(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse Word document"""
        try:
            # Ensure Office packages are installed
            if not PackageManager.ensure_packages('office'):
                raise ImportError("Office parsing packages not available")

            from docx import Document

            result = {
                "type": "word",
                "file_path": file_path,
                "metadata": {},
                "content": {
                    "text": "",
                    "paragraphs": [],
                    "tables": [],
                    "styles": []
                }
            }

            # Parse document
            doc = Document(file_path)

            # Extract metadata
            core_props = doc.core_properties
            result["metadata"] = {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
                "subject": core_props.subject or "",
                "keywords": core_props.keywords or ""
            }

            # Extract content
            full_text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    result["content"]["paragraphs"].append({
                        "text": paragraph.text,
                        "style": paragraph.style.name if paragraph.style else "Normal"
                    })
                    full_text.append(paragraph.text)

            # Extract tables
            for i, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)

                result["content"]["tables"].append({
                    "table_index": i,
                    "data": table_data
                })

            result["content"]["text"] = "\n\n".join(full_text)

            return result

        except Exception as e:
            logger.error(f"Word parsing failed: {e}")
            return {"error": str(e), "type": "word", "file_path": file_path}

    def parse_excel(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse Excel spreadsheet"""
        try:
            # Ensure Office packages are installed
            if not PackageManager.ensure_packages('office'):
                raise ImportError("Office parsing packages not available")

            import openpyxl

            result = {
                "type": "excel",
                "file_path": file_path,
                "metadata": {},
                "content": {
                    "worksheets": [],
                    "summary": ""
                }
            }

            # Load workbook
            workbook = openpyxl.load_workbook(file_path, read_only=True)

            # Extract metadata
            result["metadata"] = {
                "worksheet_count": len(workbook.worksheets),
                "worksheet_names": [ws.title for ws in workbook.worksheets]
            }

            # Extract content from each worksheet
            summary_parts = []
            for worksheet in workbook.worksheets:
                ws_data = {
                    "name": worksheet.title,
                    "max_row": worksheet.max_row,
                    "max_column": worksheet.max_column,
                    "data": []
                }

                # Read data (limit to reasonable size)
                max_rows = min(worksheet.max_row, 1000)
                max_cols = min(worksheet.max_column, 50)

                for row in worksheet.iter_rows(min_row=1, max_row=max_rows,
                                             min_col=1, max_col=max_cols, values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(cell.strip() for cell in row_data):  # Skip empty rows
                        ws_data["data"].append(row_data)

                result["content"]["worksheets"].append(ws_data)
                summary_parts.append(f"Worksheet '{worksheet.title}': {len(ws_data['data'])} rows of data")

            result["content"]["summary"] = "\n".join(summary_parts)

            return result

        except Exception as e:
            logger.error(f"Excel parsing failed: {e}")
            return {"error": str(e), "type": "excel", "file_path": file_path}

    def parse_powerpoint(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse PowerPoint presentation"""
        try:
            # Ensure Office packages are installed
            if not PackageManager.ensure_packages('office'):
                raise ImportError("Office parsing packages not available")

            from pptx import Presentation

            result = {
                "type": "powerpoint",
                "file_path": file_path,
                "metadata": {},
                "content": {
                    "text": "",
                    "slides": [],
                    "notes": []
                }
            }

            # Load presentation
            prs = Presentation(file_path)

            # Extract metadata
            result["metadata"] = {
                "slide_count": len(prs.slides),
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height
            }

            # Extract content
            full_text = []
            for i, slide in enumerate(prs.slides):
                slide_data = {
                    "slide_number": i + 1,
                    "title": "",
                    "content": [],
                    "notes": ""
                }

                # Extract slide content
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if shape.placeholder_format and shape.placeholder_format.type == 1:  # Title
                            slide_data["title"] = shape.text
                        else:
                            slide_data["content"].append(shape.text)
                        full_text.append(shape.text)

                # Extract notes
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    slide_data["notes"] = notes_text
                    if notes_text.strip():
                        full_text.append(f"Notes: {notes_text}")

                result["content"]["slides"].append(slide_data)

            result["content"]["text"] = "\n\n".join(full_text)

            return result

        except Exception as e:
            logger.error(f"PowerPoint parsing failed: {e}")
            return {"error": str(e), "type": "powerpoint", "file_path": file_path}

    def parse_image(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse image using advanced OCR engines"""
        try:
            result = {
                "type": "image",
                "file_path": file_path,
                "metadata": {},
                "content": {
                    "text": "",
                    "confidence": 0,
                    "words": [],
                    "lines": [],
                    "ocr_results": []
                }
            }

            # Load image metadata
            try:
                from PIL import Image
                with Image.open(file_path) as img:
                    result["metadata"] = {
                        "width": img.width,
                        "height": img.height,
                        "mode": img.mode,
                        "format": img.format
                    }
            except Exception as e:
                logger.warning(f"Could not load image metadata: {e}")

            # Use advanced OCR engines if available
            if OCR_AVAILABLE and ocr_manager:
                # Get OCR engine preference
                ocr_engine = kwargs.get('ocr_engine', 'auto')
                use_fallback = kwargs.get('use_fallback', True)

                if ocr_engine == 'auto' or use_fallback:
                    # Use fallback mechanism
                    engines = kwargs.get('engines', None)
                    ocr_result = ocr_manager.recognize_with_fallback(file_path, engines)
                else:
                    # Use specific engine
                    ocr_result = ocr_manager.recognize(file_path, ocr_engine)

                if ocr_result.success:
                    result["content"]["text"] = ocr_result.text
                    result["content"]["confidence"] = ocr_result.confidence
                    result["content"]["words"] = ocr_result.words
                    result["content"]["ocr_results"].append(ocr_result.to_dict())
                    logger.info(f"OCR successful with {ocr_result.provider}")
                else:
                    logger.warning(f"Advanced OCR failed: {ocr_result.error}")
                    # Fallback to pytesseract
                    return self._fallback_ocr(file_path, result, **kwargs)
            else:
                # Fallback to pytesseract
                return self._fallback_ocr(file_path, result, **kwargs)

            return result

        except Exception as e:
            logger.error(f"Image OCR failed: {e}")
            return {"error": str(e), "type": "image", "file_path": file_path}

    def _fallback_ocr(self, file_path: str, result: Dict[str, Any], **kwargs) -> Dict[str, Any]:
        """Fallback OCR using pytesseract"""
        try:
            # Ensure OCR packages are installed
            if not PackageManager.ensure_packages('ocr'):
                raise ImportError("OCR parsing packages not available")

            import pytesseract
            from PIL import Image

            logger.info("Using pytesseract fallback OCR")

            # Load and analyze image
            with Image.open(file_path) as img:
                # OCR extraction
                ocr_confidence = kwargs.get('ocr_confidence', 30)
                ocr_data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)

                # Filter confident text
                confident_text = []
                confidences = []
                words = []

                for i, confidence in enumerate(ocr_data['conf']):
                    if int(confidence) > ocr_confidence:
                        text = ocr_data['text'][i].strip()
                        if text:
                            confident_text.append(text)
                            confidences.append(int(confidence))
                            words.append({
                                "text": text,
                                "confidence": int(confidence),
                                "bbox": {
                                    "left": ocr_data['left'][i],
                                    "top": ocr_data['top'][i],
                                    "width": ocr_data['width'][i],
                                    "height": ocr_data['height'][i]
                                }
                            })

                result["content"]["text"] = " ".join(confident_text)
                result["content"]["confidence"] = sum(confidences) / len(confidences) if confidences else 0
                result["content"]["words"] = words
                result["content"]["ocr_results"].append({
                    "provider": "pytesseract",
                    "success": True,
                    "text": result["content"]["text"],
                    "confidence": result["content"]["confidence"]
                })

            return result

        except Exception as e:
            logger.error(f"Fallback OCR failed: {e}")
            result["content"]["ocr_results"].append({
                "provider": "pytesseract",
                "success": False,
                "error": str(e)
            })
            return result

    def parse_text(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse text file"""
        try:
            result = {
                "type": "text",
                "file_path": file_path,
                "metadata": {},
                "content": {
                    "text": "",
                    "lines": 0,
                    "encoding": "utf-8"
                }
            }

            # Try different encodings
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
            content = None

            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    result["content"]["encoding"] = encoding
                    break
                except UnicodeDecodeError:
                    continue

            if content is None:
                # Fallback to binary read
                with open(file_path, 'rb') as f:
                    content = f.read().decode('utf-8', errors='ignore')
                result["content"]["encoding"] = "utf-8 (with errors ignored)"

            result["content"]["text"] = content
            result["content"]["lines"] = len(content.splitlines())
            result["metadata"]["file_size"] = os.path.getsize(file_path)

            return result

        except Exception as e:
            logger.error(f"Text parsing failed: {e}")
            return {"error": str(e), "type": "text", "file_path": file_path}

    def parse_json(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse JSON file"""
        try:
            result = {
                "type": "json",
                "file_path": file_path,
                "metadata": {},
                "content": {}
            }

            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)

            result["content"] = data
            result["metadata"]["file_size"] = os.path.getsize(file_path)

            return result

        except Exception as e:
            logger.error(f"JSON parsing failed: {e}")
            return {"error": str(e), "type": "json", "file_path": file_path}

    def parse_xml(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse XML file"""
        try:
            # Ensure XML packages are installed
            if not PackageManager.ensure_packages('xml'):
                raise ImportError("XML parsing packages not available")

            import xml.etree.ElementTree as ET

            result = {
                "type": "xml",
                "file_path": file_path,
                "metadata": {},
                "content": {
                    "text": "",
                    "structure": {}
                }
            }

            # Parse XML
            tree = ET.parse(file_path)
            root = tree.getroot()

            def xml_to_dict(element):
                result = {}
                if element.text and element.text.strip():
                    result['text'] = element.text.strip()
                if element.attrib:
                    result['attributes'] = element.attrib

                children = {}
                for child in element:
                    child_data = xml_to_dict(child)
                    if child.tag in children:
                        if not isinstance(children[child.tag], list):
                            children[child.tag] = [children[child.tag]]
                        children[child.tag].append(child_data)
                    else:
                        children[child.tag] = child_data

                if children:
                    result['children'] = children

                return result

            result["content"]["structure"] = {root.tag: xml_to_dict(root)}

            # Extract all text content
            def extract_text(element):
                text_parts = []
                if element.text and element.text.strip():
                    text_parts.append(element.text.strip())
                for child in element:
                    text_parts.extend(extract_text(child))
                return text_parts

            all_text = extract_text(root)
            result["content"]["text"] = "\n".join(all_text)
            result["metadata"]["file_size"] = os.path.getsize(file_path)

            return result

        except Exception as e:
            logger.error(f"XML parsing failed: {e}")
            return {"error": str(e), "type": "xml", "file_path": file_path}

    def parse_file(self, file_path: str, output_format: str = "json", **kwargs) -> Dict[str, Any]:
        """Parse any supported file format"""
        try:
            # Validate file exists
            if not os.path.exists(file_path):
                return {"error": f"File not found: {file_path}"}

            # Detect file type
            file_type = self.detect_file_type(file_path)
            if not file_type:
                return {"error": f"Unsupported file format: {file_path}"}

            # Parse file
            parser_func = self.supported_formats[file_type]
            result = parser_func(file_path, **kwargs)

            # Add general metadata
            if "metadata" not in result:
                result["metadata"] = {}

            result["metadata"].update({
                "file_name": os.path.basename(file_path),
                "file_size": os.path.getsize(file_path),
                "detected_type": file_type,
                "parser_version": "1.0"
            })

            return result

        except Exception as e:
            logger.error(f"File parsing failed: {e}")
            return {"error": str(e), "file_path": file_path}

# Initialize parser and converter
parser = DocumentParser()
converter = DocumentConverter()

# Initialize OCR queue system
if OCR_AVAILABLE and ocr_queue:
    ocr_queue.start_processing()
    logger.info("OCR queue system started")

# Initialize packages in background
PackageManager.initialize_all_packages()

@mcp.tool()
def initialize_packages(
    format_types: Optional[List[str]] = None,
    force_reinstall: bool = False
) -> Dict[str, Any]:
    """
    Initialize and install required packages for document processing.

    Args:
        format_types: List of format types to initialize (pdf, xmind, office, etc.)
                     If None, initializes all packages
        force_reinstall: Force reinstallation of packages

    Returns:
        Dictionary containing initialization results
    """
    try:
        logger.info("Starting package initialization...")

        if force_reinstall:
            # Clear installation cache
            PackageManager._installation_status.clear()

        # Determine which packages to install
        if format_types is None:
            format_types = list(PackageManager.PACKAGE_MAPPING.keys())

        results = {
            "requested_formats": format_types,
            "installation_results": {},
            "successful_formats": [],
            "failed_formats": []
        }

        # Install packages for each format type
        for format_type in format_types:
            if format_type in PackageManager.PACKAGE_MAPPING:
                logger.info(f"Initializing packages for {format_type}...")
                success = PackageManager.ensure_packages(format_type)

                results["installation_results"][format_type] = {
                    "success": success,
                    "packages": PackageManager.PACKAGE_MAPPING[format_type]
                }

                if success:
                    results["successful_formats"].append(format_type)
                else:
                    results["failed_formats"].append(format_type)
            else:
                results["installation_results"][format_type] = {
                    "success": False,
                    "error": "Unknown format type"
                }
                results["failed_formats"].append(format_type)

        logger.info(f"Package initialization completed. Success: {len(results['successful_formats'])}, Failed: {len(results['failed_formats'])}")
        return results

    except Exception as e:
        error_msg = f"Package initialization failed: {str(e)}"
        logger.error(error_msg)
        return {"error": error_msg}

@mcp.tool()
def parse_document(
    file_path: str,
    output_format: str = "json",
    extract_tables: bool = True,
    extract_images: bool = False,
    ocr_confidence: int = 30
) -> Dict[str, Any]:
    """
    Parse various document formats and extract AI-readable content.

    Supported formats:
    - PDF files (text extraction and OCR)
    - XMind files (mind maps)
    - Word documents (.docx, .doc)
    - Excel spreadsheets (.xlsx, .xls)
    - PowerPoint presentations (.pptx, .ppt)
    - Images (OCR extraction)
    - Text files (.txt, .md, .json, .xml)

    Args:
        file_path: Path to the document file
        output_format: Output format (json, xml, markdown)
        extract_tables: Whether to extract tables (for supported formats)
        extract_images: Whether to extract images (for supported formats)
        ocr_confidence: Minimum OCR confidence threshold (0-100)

    Returns:
        Dictionary containing parsed content, metadata, and structure
    """
    try:
        # Normalize path
        file_path = os.path.abspath(file_path)

        logger.info(f"Parsing document: {file_path}")
        logger.info(f"Output format: {output_format}")

        # Parse document
        result = parser.parse_file(
            file_path=file_path,
            output_format=output_format,
            extract_tables=extract_tables,
            extract_images=extract_images,
            ocr_confidence=ocr_confidence
        )

        # Format output based on requested format
        if output_format.lower() == "xml" and "error" not in result:
            try:
                if PackageManager.ensure_packages('xml'):
                    from dicttoxml import dicttoxml
                    xml_content = dicttoxml(result, custom_root='document', attr_type=False)
                    result["xml_output"] = xml_content.decode('utf-8')
            except Exception as e:
                logger.warning(f"XML conversion failed: {e}")

        elif output_format.lower() == "markdown" and "error" not in result:
            try:
                if PackageManager.ensure_packages('markdown'):
                    # Convert to markdown format
                    md_parts = []

                    # Add title
                    if "metadata" in result and "file_name" in result["metadata"]:
                        md_parts.append(f"# {result['metadata']['file_name']}")

                    # Add metadata
                    if "metadata" in result:
                        md_parts.append("## Metadata")
                        for key, value in result["metadata"].items():
                            md_parts.append(f"- **{key}**: {value}")

                    # Add content
                    if "content" in result:
                        md_parts.append("## Content")
                        if "text" in result["content"]:
                            md_parts.append(result["content"]["text"])

                    result["markdown_output"] = "\n\n".join(md_parts)
            except Exception as e:
                logger.warning(f"Markdown conversion failed: {e}")

        logger.info(f"Document parsing completed: {result.get('type', 'unknown')} format")
        return result

    except Exception as e:
        error_msg = f"Document parsing failed: {str(e)}"
        logger.error(error_msg)
        return {"error": error_msg, "file_path": file_path}

@mcp.tool()
def list_supported_formats() -> Dict[str, Any]:
    """
    List all supported document formats and their capabilities.

    Returns:
        Dictionary containing supported formats and their features
    """
    return {
        "supported_formats": {
            "PDF": {
                "extensions": [".pdf"],
                "features": ["text_extraction", "table_extraction", "metadata"],
                "description": "Portable Document Format - text and table extraction"
            },
            "XMind": {
                "extensions": [".xmind"],
                "features": ["structure_extraction", "text_extraction"],
                "description": "Mind mapping files - structure and content extraction"
            },
            "Word": {
                "extensions": [".docx", ".doc"],
                "features": ["text_extraction", "table_extraction", "metadata", "style_preservation"],
                "description": "Microsoft Word documents"
            },
            "Excel": {
                "extensions": [".xlsx", ".xls"],
                "features": ["data_extraction", "worksheet_separation", "metadata"],
                "description": "Microsoft Excel spreadsheets"
            },
            "PowerPoint": {
                "extensions": [".pptx", ".ppt"],
                "features": ["text_extraction", "slide_separation", "notes_extraction"],
                "description": "Microsoft PowerPoint presentations"
            },
            "Images": {
                "extensions": [".jpg", ".jpeg", ".png", ".bmp", ".tiff"],
                "features": ["ocr_extraction", "confidence_filtering"],
                "description": "Image files with OCR text extraction"
            },
            "Text": {
                "extensions": [".txt", ".md"],
                "features": ["encoding_detection", "line_counting"],
                "description": "Plain text and Markdown files"
            },
            "Structured": {
                "extensions": [".json", ".xml"],
                "features": ["structure_preservation", "data_extraction"],
                "description": "JSON and XML structured data files"
            }
        },
        "output_formats": ["json", "xml", "markdown"],
        "ai_optimizations": [
            "Smart content extraction",
            "Structure preservation",
            "Metadata extraction",
            "Error handling",
            "Multi-format support"
        ]
    }

@mcp.tool()
def convert_document(
    input_data: str,
    from_format: str,
    to_format: str,
    output_path: Optional[str] = None,
    title: str = "Document",
    page_size: str = "A4",
    encoding: str = "UTF-8",
    ignore_links: bool = False,
    ignore_images: bool = False
) -> Dict[str, Any]:
    """
    Convert between different document formats.

    Supported conversions:
    - HTML to Markdown, PDF, Word (DOCX)
    - Markdown to HTML, PDF, Word (DOCX)
    - Text to HTML
    - JSON to Excel
    - Excel to JSON (requires file path as input_data)

    Args:
        input_data: Input content (string) or file path (for Excel to JSON)
        from_format: Source format (html, markdown, text, json, excel)
        to_format: Target format (markdown, html, pdf, docx, excel, json)
        output_path: Output file path (required for file formats like PDF, DOCX, Excel)
        title: Document title (for HTML generation)
        page_size: PDF page size (A4, Letter, etc.)
        encoding: Text encoding
        ignore_links: Ignore links in HTML to Markdown conversion
        ignore_images: Ignore images in HTML to Markdown conversion

    Returns:
        Dictionary containing conversion result or converted content
    """
    try:
        logger.info(f"Converting from {from_format} to {to_format}")

        # Special handling for Excel to JSON (requires file path)
        if from_format == 'excel' and to_format == 'json':
            if not os.path.exists(input_data):
                return {"error": f"Excel file not found: {input_data}"}
            result_data = converter.excel_to_json(input_data)
            return {
                "success": True,
                "from_format": from_format,
                "to_format": to_format,
                "data": result_data
            }

        # For JSON to Excel, parse input if it's a string
        if from_format == 'json':
            if isinstance(input_data, str):
                try:
                    input_data = json.loads(input_data)
                except json.JSONDecodeError as e:
                    return {"error": f"Invalid JSON data: {e}"}

        # Perform conversion
        result = converter.convert_format(
            input_data=input_data,
            from_format=from_format,
            to_format=to_format,
            output_path=output_path,
            title=title,
            page_size=page_size,
            encoding=encoding,
            ignore_links=ignore_links,
            ignore_images=ignore_images
        )

        if output_path and isinstance(result, bool):
            # File was created
            if result:
                return {
                    "success": True,
                    "from_format": from_format,
                    "to_format": to_format,
                    "output_path": output_path,
                    "file_created": True
                }
            else:
                return {"error": "Conversion failed"}
        else:
            # String result
            return {
                "success": True,
                "from_format": from_format,
                "to_format": to_format,
                "content": result
            }

    except Exception as e:
        error_msg = f"Document conversion failed: {str(e)}"
        logger.error(error_msg)
        return {"error": error_msg}

@mcp.tool()
def write_document(
    content: str,
    output_path: str,
    format_type: str = "auto",
    metadata: Optional[Dict[str, Any]] = None,
    encoding: str = "utf-8"
) -> Dict[str, Any]:
    """
    Write content to a document file in various formats.

    Args:
        content: Content to write
        output_path: Output file path
        format_type: Document format (auto, text, json, html, markdown, xml)
        metadata: Optional metadata to include
        encoding: Text encoding

    Returns:
        Dictionary containing write operation result
    """
    try:
        logger.info(f"Writing document to: {output_path}")

        # Auto-detect format from file extension
        if format_type == "auto":
            ext = Path(output_path).suffix.lower()
            format_mapping = {
                '.txt': 'text',
                '.md': 'markdown',
                '.html': 'html',
                '.htm': 'html',
                '.json': 'json',
                '.xml': 'xml'
            }
            format_type = format_mapping.get(ext, 'text')

        # Ensure output directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        if format_type == 'json':
            # Write as JSON
            try:
                data = json.loads(content) if isinstance(content, str) else content
                if metadata:
                    data = {"metadata": metadata, "content": data}
                with open(output_path, 'w', encoding=encoding) as f:
                    json.dump(data, f, ensure_ascii=False, indent=2)
            except json.JSONDecodeError:
                # Content is not JSON, wrap it
                data = {"content": content}
                if metadata:
                    data["metadata"] = metadata
                with open(output_path, 'w', encoding=encoding) as f:
                    json.dump(data, f, ensure_ascii=False, indent=2)

        elif format_type == 'html':
            # Write as HTML
            if not content.strip().startswith('<!DOCTYPE html>'):
                # Convert plain text to HTML
                html_content = converter.text_to_html(
                    content,
                    title=metadata.get('title', 'Document') if metadata else 'Document'
                )
                content = html_content

            with open(output_path, 'w', encoding=encoding) as f:
                f.write(content)

        elif format_type == 'xml':
            # Write as XML
            if not content.strip().startswith('<?xml'):
                # Wrap content in XML
                xml_content = f'<?xml version="1.0" encoding="{encoding}"?>\n<document>\n<content>{content}</content>\n</document>'
                content = xml_content

            with open(output_path, 'w', encoding=encoding) as f:
                f.write(content)

        else:
            # Write as plain text
            with open(output_path, 'w', encoding=encoding) as f:
                f.write(content)

        file_size = os.path.getsize(output_path)

        return {
            "success": True,
            "output_path": output_path,
            "format_type": format_type,
            "file_size": file_size,
            "encoding": encoding
        }

    except Exception as e:
        error_msg = f"Document write failed: {str(e)}"
        logger.error(error_msg)
        return {"error": error_msg}

@mcp.tool()
def batch_convert(
    input_directory: str,
    output_directory: str,
    from_format: str,
    to_format: str,
    file_pattern: str = "*",
    recursive: bool = False
) -> Dict[str, Any]:
    """
    Batch convert multiple documents from one format to another.

    Args:
        input_directory: Directory containing input files
        output_directory: Directory for output files
        from_format: Source format
        to_format: Target format
        file_pattern: File pattern to match (e.g., "*.html")
        recursive: Search subdirectories recursively

    Returns:
        Dictionary containing batch conversion results
    """
    try:
        from pathlib import Path
        import glob

        logger.info(f"Batch converting {from_format} to {to_format}")
        logger.info(f"Input: {input_directory}, Output: {output_directory}")

        # Ensure output directory exists
        os.makedirs(output_directory, exist_ok=True)

        # Find files
        search_pattern = os.path.join(input_directory, "**" if recursive else "", file_pattern)
        files = glob.glob(search_pattern, recursive=recursive)

        results = {
            "total_files": len(files),
            "successful": 0,
            "failed": 0,
            "details": []
        }

        for file_path in files:
            try:
                # Read input file
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()

                # Generate output filename
                input_name = Path(file_path).stem
                output_ext = {
                    'html': '.html',
                    'markdown': '.md',
                    'pdf': '.pdf',
                    'docx': '.docx',
                    'text': '.txt',
                    'json': '.json',
                    'excel': '.xlsx'
                }.get(to_format, '.txt')

                output_path = os.path.join(output_directory, f"{input_name}{output_ext}")

                # Convert
                result = converter.convert_format(
                    input_data=content,
                    from_format=from_format,
                    to_format=to_format,
                    output_path=output_path if to_format in ['pdf', 'docx', 'excel'] else None
                )

                if to_format not in ['pdf', 'docx', 'excel'] and isinstance(result, str):
                    # Save string result
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(result)

                results["successful"] += 1
                results["details"].append({
                    "input_file": file_path,
                    "output_file": output_path,
                    "status": "success"
                })

            except Exception as e:
                results["failed"] += 1
                results["details"].append({
                    "input_file": file_path,
                    "error": str(e),
                    "status": "failed"
                })
                logger.error(f"Failed to convert {file_path}: {e}")

        return results

    except Exception as e:
        error_msg = f"Batch conversion failed: {str(e)}"
        logger.error(error_msg)
        return {"error": error_msg}

@mcp.tool()
def list_conversion_formats() -> Dict[str, Any]:
    """
    List all supported format conversions.

    Returns:
        Dictionary containing all supported conversion mappings
    """
    return {
        "supported_conversions": {
            "html": ["markdown", "pdf", "docx"],
            "markdown": ["html", "pdf", "docx"],
            "text": ["html"],
            "json": ["excel"],
            "excel": ["json"]
        },
        "file_output_formats": ["pdf", "docx", "excel"],
        "string_output_formats": ["html", "markdown", "json"],
        "write_formats": ["text", "html", "markdown", "json", "xml"],
        "batch_supported": True,
        "ai_optimizations": [
            "Smart content extraction",
            "Format-aware conversion",
            "Metadata preservation",
            "Batch processing",
            "Error handling"
        ]
    }

@mcp.tool()
def ocr_recognize(
    image_path: str,
    ocr_engine: str = "auto",
    language: str = "chs",
    use_fallback: bool = True,
    confidence_threshold: int = 30,
    tencent_secret_id: Optional[str] = None,
    tencent_secret_key: Optional[str] = None
) -> Dict[str, Any]:
    """
    Perform OCR recognition on image using multiple OCR engines.

    Supported OCR engines:
    - free: Free OCR.space API (default)
    - tencent: Tencent Cloud OCR (requires credentials)
    - auto: Automatic fallback through available engines

    Args:
        image_path: Path to the image file
        ocr_engine: OCR engine to use (free, tencent, auto)
        language: Language for OCR (chs=Chinese, eng=English, auto=Auto-detect)
        use_fallback: Use fallback to other engines if primary fails
        confidence_threshold: Minimum confidence threshold for text extraction (0-100)
        tencent_secret_id: Tencent Cloud Secret ID (optional)
        tencent_secret_key: Tencent Cloud Secret Key (optional)

    Returns:
        Dictionary containing OCR results with text, confidence, and word details
    """
    try:
        # Normalize file path for Chinese characters and Windows paths
        normalized_path = _normalize_file_path(image_path)
        logger.info(f"OCR recognition request: {normalized_path}, engine: {ocr_engine}")

        # Validate file exists and is accessible
        if not _validate_image_file(normalized_path):
            return {
                "success": False,
                "error": f"Image file not found or not accessible: {image_path}",
                "details": "Please check the file path and ensure the file exists"
            }

        # Check if OCR is available
        if not OCR_AVAILABLE or not ocr_manager:
            return {"error": "OCR engines not available"}

        # Set Tencent credentials if provided
        if tencent_secret_id and tencent_secret_key:
            ocr_manager.set_tencent_credentials(tencent_secret_id, tencent_secret_key)

        # Prepare OCR parameters
        ocr_params = {
            'language': language,
            'confidence_threshold': confidence_threshold
        }

        # Perform OCR with enhanced fallback logic
        if ocr_engine == "auto" or use_fallback:
            available_engines = ocr_manager.get_available_engines()
            logger.info(f"Available OCR engines: {available_engines}")

            # Use normalized path for OCR processing
            ocr_result = ocr_manager.recognize_with_fallback(normalized_path, available_engines, **ocr_params)

            # If fallback also fails, try individual engines with more detailed error reporting
            if not ocr_result.success and len(available_engines) > 1:
                logger.warning(f"Fallback OCR failed: {ocr_result.error}. Trying individual engines...")

                for engine in available_engines:
                    try:
                        logger.info(f"Attempting OCR with {engine} engine...")
                        single_result = ocr_manager.recognize(normalized_path, engine, **ocr_params)

                        if single_result.success:
                            logger.info(f"OCR succeeded with {engine} engine")
                            ocr_result = single_result
                            break
                        else:
                            logger.warning(f"OCR failed with {engine}: {single_result.error}")
                    except Exception as engine_error:
                        logger.error(f"Engine {engine} crashed: {engine_error}")
                        continue
        else:
            ocr_result = ocr_manager.recognize(normalized_path, ocr_engine, **ocr_params)

        # Format response
        response = {
            "success": ocr_result.success,
            "provider": ocr_result.provider,
            "processing_time": ocr_result.processing_time,
            "image_path": image_path
        }

        if ocr_result.success:
            response.update({
                "text": ocr_result.text,
                "confidence": ocr_result.confidence,
                "word_count": len(ocr_result.words),
                "words": ocr_result.words[:50],  # Limit to first 50 words
                "has_more_words": len(ocr_result.words) > 50
            })
        else:
            response["error"] = ocr_result.error

        logger.info(f"OCR completed: {response['success']}, provider: {response['provider']}")
        return response

    except Exception as e:
        error_msg = f"OCR recognition failed: {str(e)}"
        logger.error(error_msg)
        return {"error": error_msg}

@mcp.tool()
def configure_ocr(
    free_ocr_api_key: Optional[str] = None,
    tencent_secret_id: Optional[str] = None,
    tencent_secret_key: Optional[str] = None,
    tencent_region: str = "ap-beijing",
    test_connection: bool = False
) -> Dict[str, Any]:
    """
    Configure OCR engines with credentials and test connectivity.

    Args:
        free_ocr_api_key: Free OCR (OCR.space) API key
        tencent_secret_id: Tencent Cloud Secret ID
        tencent_secret_key: Tencent Cloud Secret Key
        tencent_region: Tencent Cloud region
        test_connection: Test OCR engine connectivity

    Returns:
        Dictionary containing configuration status and available engines
    """
    try:
        logger.info("Configuring OCR engines...")

        if not OCR_AVAILABLE or not ocr_manager:
            return {"error": "OCR engines not available"}

        result = {
            "success": True,
            "configured_engines": [],
            "available_engines": [],
            "test_results": {}
        }

        # Configure Free OCR
        if free_ocr_api_key:
            try:
                ocr_manager.set_free_ocr_key(free_ocr_api_key)
                result["configured_engines"].append("free")
                logger.info("Free OCR API key configured")
            except Exception as e:
                result["test_results"]["free"] = {"error": str(e)}
                logger.error(f"Free OCR configuration failed: {e}")

        # Configure Tencent Cloud
        if tencent_secret_id and tencent_secret_key:
            try:
                ocr_manager.set_tencent_credentials(tencent_secret_id, tencent_secret_key, tencent_region)
                result["configured_engines"].append("tencent")
                logger.info("Tencent Cloud OCR configured")
            except Exception as e:
                result["test_results"]["tencent"] = {"error": str(e)}
                logger.error(f"Tencent OCR configuration failed: {e}")

        # Get available engines
        result["available_engines"] = ocr_manager.get_available_engines()

        # Test connectivity if requested
        if test_connection:
            # Create a simple test image
            test_image_path = None
            try:
                import tempfile
                from PIL import Image, ImageDraw, ImageFont

                # Create test image
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                    test_image_path = tmp.name

                # Create simple test image with text
                img = Image.new('RGB', (200, 100), color='white')
                draw = ImageDraw.Draw(img)
                draw.text((10, 10), "Test OCR", fill='black')
                img.save(test_image_path)

                # Test each available engine
                for engine in result["available_engines"]:
                    try:
                        test_result = ocr_manager.recognize(test_image_path, engine)
                        result["test_results"][engine] = {
                            "success": test_result.success,
                            "text": test_result.text if test_result.success else None,
                            "error": test_result.error if not test_result.success else None,
                            "processing_time": test_result.processing_time
                        }
                    except Exception as e:
                        result["test_results"][engine] = {"error": str(e)}

            except Exception as e:
                result["test_results"]["general"] = {"error": f"Test setup failed: {e}"}
            finally:
                # Clean up test image
                if test_image_path and os.path.exists(test_image_path):
                    try:
                        os.unlink(test_image_path)
                    except:
                        pass

        logger.info(f"OCR configuration completed: {len(result['available_engines'])} engines available")
        return result

    except Exception as e:
        error_msg = f"OCR configuration failed: {str(e)}"
        logger.error(error_msg)
        return {"error": error_msg}

@mcp.tool()
def list_ocr_engines() -> Dict[str, Any]:
    """
    List all available OCR engines and their capabilities.

    Returns:
        Dictionary containing OCR engine information and capabilities
    """
    try:
        if not OCR_AVAILABLE or not ocr_manager:
            return {
                "available": False,
                "error": "OCR engines not available",
                "engines": {}
            }

        available_engines = ocr_manager.get_available_engines()

        engines_info = {
            "free": {
                "name": "Free OCR",
                "provider": "OCR.space",
                "cost": "Free",
                "limits": "1MB file size, rate limited",
                "languages": ["English", "Chinese", "Auto-detect"],
                "features": ["Text extraction", "Word coordinates", "Confidence scores"],
                "setup_required": False
            },
            "tencent": {
                "name": "Tencent Cloud OCR",
                "provider": "Tencent Cloud",
                "cost": "Paid (1000 free calls/month)",
                "limits": "High accuracy, enterprise grade",
                "languages": ["Chinese", "English", "Multi-language"],
                "features": ["High accuracy", "Multi-language", "Special document types"],
                "setup_required": True,
                "credentials": ["TENCENT_SECRET_ID", "TENCENT_SECRET_KEY"]
            }
        }

        return {
            "available": True,
            "available_engines": available_engines,
            "total_engines": len(engines_info),
            "engines": {name: info for name, info in engines_info.items()},
            "recommendations": {
                "for_chinese": "tencent",
                "for_english": "free",
                "for_production": "tencent",
                "for_testing": "free"
            }
        }

    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def smart_ocr_recognize(
    file_paths: Union[str, List[str], List[List[str]]],
    priority: str = "normal",
    content_type: str = "general",
    batch_processing: bool = False,
    wait_for_completion: bool = True,
    timeout: int = 300
) -> Dict[str, Any]:
    """
    Intelligent OCR recognition with smart resource management.

    Supports single files, batches, and 2D queues with automatic optimization:
    - Images are automatically compressed to meet engine limits (1MB, 720p max)
    - PDFs are split into chunks (3 pages max for free OCR)
    - Smart engine selection based on quotas and content type
    - Queue processing for large batches

    Args:
        file_paths: Single file, list of files, or 2D array [[group1], [group2]]
        priority: Processing priority ("high", "normal", "low")
        content_type: Content type hint ("chinese_text", "english_text", "mixed_content", "general")
        batch_processing: Force batch processing mode
        wait_for_completion: Wait for processing to complete
        timeout: Maximum wait time in seconds

    Returns:
        Dictionary containing OCR results or processing status
    """
    try:
        if not OCR_AVAILABLE or not ocr_queue:
            return {"error": "OCR system not available"}

        logger.info(f"Smart OCR request: {type(file_paths)} with {len(file_paths) if isinstance(file_paths, list) else 1} items")

        # Parse priority
        priority_map = {
            "high": TaskPriority.HIGH,
            "normal": TaskPriority.NORMAL,
            "low": TaskPriority.LOW
        }
        task_priority = priority_map.get(priority.lower(), TaskPriority.NORMAL)

        # Normalize all inputs to 2D array format for compatibility
        normalized_2d_queue = []

        if isinstance(file_paths, str):
            # Single file -> [[file]]
            # Normalize file path for Chinese characters and Windows paths
            normalized_path = _normalize_file_path(file_paths)
            if not _validate_image_file(normalized_path):
                return {
                    "success": False,
                    "error": f"Image file not found or not accessible: {file_paths}",
                    "details": "Please check the file path and ensure the file exists"
                }
            normalized_2d_queue = [[normalized_path]]
            logger.info("Single file converted to 2D queue format")

        elif isinstance(file_paths, list) and len(file_paths) > 0:
            # Check if it's already a 2D array
            is_2d = isinstance(file_paths[0], list)

            if is_2d:
                # Already 2D format - normalize all paths in nested structure
                normalized_2d_queue = []
                for group in file_paths:
                    normalized_group = []
                    for file_path in group:
                        normalized_path = _normalize_file_path(file_path)
                        if _validate_image_file(normalized_path):
                            normalized_group.append(normalized_path)
                        else:
                            logger.warning(f"Skipping invalid file: {file_path}")
                    if normalized_group:  # Only add non-empty groups
                        normalized_2d_queue.append(normalized_group)
                logger.info("Input 2D queue normalized")
            else:
                # 1D list -> [file_paths] with normalization
                normalized_group = []
                for file_path in file_paths:
                    normalized_path = _normalize_file_path(file_path)
                    if _validate_image_file(normalized_path):
                        normalized_group.append(normalized_path)
                    else:
                        logger.warning(f"Skipping invalid file: {file_path}")

                if normalized_group:
                    normalized_2d_queue = [normalized_group]
                    logger.info("1D list converted to 2D queue format")
                else:
                    return {
                        "success": False,
                        "error": "No valid image files found in the provided list",
                        "details": "Please check that all file paths are correct and accessible"
                    }
        else:
            return {"error": "Invalid file_paths format"}

        # Process as 2D queue (all requests are now 2D arrays)
        batch_info = ocr_queue.add_2d_queue(
            queue_data=normalized_2d_queue,
            priority=task_priority
        )

        if wait_for_completion:
            result = _wait_for_batch_completion(batch_info, timeout)
            return result
        else:
            return {
                "success": True,
                "mode": "2d_queue",
                "batch_id": batch_info["batch_id"],
                "groups": batch_info["groups"],
                "total_files": batch_info["total_files"],
                "status": "queued"
            }

    except Exception as e:
        error_msg = f"Smart OCR recognition failed: {str(e)}"
        logger.error(error_msg)
        return {"error": error_msg}

@mcp.tool()
def get_ocr_task_status(task_id: str) -> Dict[str, Any]:
    """
    Get status of an OCR task.

    Args:
        task_id: Task ID returned from smart_ocr_recognize

    Returns:
        Dictionary containing task status and results
    """
    try:
        if not OCR_AVAILABLE or not ocr_queue:
            return {"error": "OCR system not available"}

        status = ocr_queue.get_task_status(task_id)
        if status:
            return {"success": True, "task_status": status}
        else:
            return {"error": f"Task {task_id} not found"}

    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def get_ocr_batch_status(group_id: str) -> Dict[str, Any]:
    """
    Get status of an OCR batch group.

    Args:
        group_id: Group ID returned from smart_ocr_recognize

    Returns:
        Dictionary containing batch status and results
    """
    try:
        if not OCR_AVAILABLE or not ocr_queue:
            return {"error": "OCR system not available"}

        # Enhanced status checking with detailed information
        status = ocr_queue.get_batch_status(group_id)
        if status:
            # Add queue system status for better debugging
            queue_stats = {
                "queue_size": ocr_queue.task_queue.qsize() if hasattr(ocr_queue.task_queue, 'qsize') else 0,
                "active_tasks": len(ocr_queue.active_tasks),
                "completed_tasks": len(ocr_queue.completed_tasks),
                "worker_count": len(ocr_queue.worker_threads)
            }

            return {
                "success": True,
                "batch_status": status,
                "queue_info": queue_stats,
                "processing_hint": "If status shows 'queued', tasks are waiting for available workers. This is normal during high load."
            }
        else:
            # Check if it might be a single task instead of batch
            task_status = ocr_queue.get_task_status(group_id) if hasattr(ocr_queue, 'get_task_status') else None
            if task_status:
                return {
                    "success": True,
                    "task_status": task_status,
                    "note": "This appears to be a single task, not a batch"
                }
            else:
                return {
                    "error": f"Batch or task {group_id} not found",
                    "suggestion": "The task may have completed and been cleaned up, or the ID might be incorrect"
                }

    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def get_ocr_system_status() -> Dict[str, Any]:
    """
    Get OCR system status and resource usage.

    Returns:
        Dictionary containing system statistics and resource limits
    """
    try:
        if not OCR_AVAILABLE or not ocr_queue:
            return {"error": "OCR system not available"}

        stats = ocr_queue.get_system_stats()

        # Add resource limits information
        limits = OCRLimits()
        resource_info = {
            "free_ocr": {
                "monthly_limit": limits.FREE_OCR["requests_per_month"],
                "file_size_limit_mb": limits.FREE_OCR["file_size_limit"] / (1024 * 1024),
                "pdf_page_limit": limits.FREE_OCR["pdf_page_limit"],
                "max_resolution": limits.FREE_OCR["max_image_resolution"]
            },
            "tencent_ocr": {
                "monthly_limit": limits.TENCENT_OCR["requests_per_month"],
                "file_size_limit_mb": limits.TENCENT_OCR["file_size_limit"] / (1024 * 1024),
                "pdf_page_limit": limits.TENCENT_OCR["pdf_page_limit"],
                "max_resolution": limits.TENCENT_OCR["max_image_resolution"]
            }
        }

        return {
            "success": True,
            "system_stats": stats,
            "resource_limits": resource_info,
            "features": {
                "smart_compression": True,
                "pdf_splitting": True,
                "queue_processing": True,
                "auto_engine_selection": True,
                "batch_processing": True,
                "2d_queue_support": True
            }
        }

    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def scan_directory_and_ocr(
    directory_path: str,
    max_depth: int = 3,
    ocr_engine: str = "auto",
    language: str = "chs",
    recursive: bool = True,
    image_extensions: Optional[List[str]] = None
) -> Dict[str, Any]:
    """
    Scan directory for images and perform OCR on all found images.

    This method scans a directory (with configurable depth) for image files,
    then performs OCR recognition on each image and returns a map of results.

    Args:
        directory_path: Directory path to scan for images
        max_depth: Maximum directory depth to scan (default: 3, set to 0 for unlimited)
        ocr_engine: OCR engine to use (free, tencent, auto) - default: auto
        language: Language for OCR (chs=Chinese, eng=English, auto=Auto-detect) - default: chs
        recursive: Whether to scan subdirectories (default: True)
        image_extensions: List of image extensions to scan (default: ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'])

    Returns:
        Dictionary containing:
        - success: bool
        - scanned_files: int (total image files found)
        - ocr_results: Dict[str, Dict] (map of filepath -> OCR result)
        - errors: List[Dict] (list of files that failed)
        - summary: Dict (statistics)
    """
    try:
        if not OCR_AVAILABLE or not ocr_manager:
            return {
                "success": False,
                "error": "OCR system not available",
                "scanned_files": 0,
                "ocr_results": {},
                "errors": []
            }

        # Normalize directory path
        dir_path = Path(directory_path)
        if not dir_path.exists():
            return {
                "success": False,
                "error": f"Directory not found: {directory_path}",
                "scanned_files": 0,
                "ocr_results": {},
                "errors": []
            }

        if not dir_path.is_dir():
            return {
                "success": False,
                "error": f"Path is not a directory: {directory_path}",
                "scanned_files": 0,
                "ocr_results": {},
                "errors": []
            }

        # Default image extensions
        if image_extensions is None:
            image_extensions = ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp', '.tiff', '.tif']

        # Normalize extensions to lowercase
        image_extensions = [ext.lower() if ext.startswith('.') else f'.{ext.lower()}' for ext in image_extensions]

        # Scan directory for images
        logger.info(f"Scanning directory: {dir_path} (max_depth={max_depth}, recursive={recursive})")
        found_images = []

        def scan_dir(current_path: Path, current_depth: int = 0):
            """Recursively scan directory with depth limit"""
            try:
                # Check depth limit (0 means unlimited)
                if max_depth > 0 and current_depth >= max_depth:
                    return

                for item in current_path.iterdir():
                    try:
                        if item.is_file():
                            # Check if file has image extension
                            if item.suffix.lower() in image_extensions:
                                found_images.append(str(item.resolve()))

                        elif item.is_dir() and recursive:
                            # Recursively scan subdirectories
                            scan_dir(item, current_depth + 1)

                    except Exception as e:
                        logger.warning(f"Error processing {item}: {e}")
                        continue

            except Exception as e:
                logger.warning(f"Error scanning directory {current_path}: {e}")

        # Start scanning from root directory
        scan_dir(dir_path, 0)

        logger.info(f"Found {len(found_images)} image files")

        if len(found_images) == 0:
            return {
                "success": True,
                "scanned_files": 0,
                "ocr_results": {},
                "errors": [],
                "summary": {
                    "total_files": 0,
                    "successful": 0,
                    "failed": 0,
                    "scan_path": str(dir_path),
                    "max_depth": max_depth,
                    "recursive": recursive
                }
            }

        # Perform OCR on each image
        ocr_results = {}
        errors = []
        successful_count = 0

        for idx, image_path in enumerate(found_images, 1):
            try:
                logger.info(f"Processing {idx}/{len(found_images)}: {image_path}")

                # Normalize file path
                normalized_path = _normalize_file_path(image_path)

                # Perform OCR
                result = ocr_manager.recognize(
                    image_path=normalized_path,
                    engine=ocr_engine,
                    language=language
                )

                if result and result.success:
                    ocr_results[image_path] = {
                        "success": True,
                        "text": result.text,
                        "confidence": result.confidence,
                        "provider": result.provider,
                        "word_count": len(result.text.split()) if result.text else 0,
                        "processing_time": result.processing_time,
                        "lines": result.lines if hasattr(result, 'lines') else [],
                        "words": result.words if hasattr(result, 'words') else []
                    }
                    successful_count += 1
                else:
                    error_msg = result.error if result else "Unknown error"
                    ocr_results[image_path] = {
                        "success": False,
                        "error": error_msg
                    }
                    errors.append({
                        "file": image_path,
                        "error": error_msg
                    })

            except Exception as e:
                error_msg = str(e)
                logger.error(f"Error processing {image_path}: {error_msg}")
                ocr_results[image_path] = {
                    "success": False,
                    "error": error_msg
                }
                errors.append({
                    "file": image_path,
                    "error": error_msg
                })

        # Generate summary
        summary = {
            "total_files": len(found_images),
            "successful": successful_count,
            "failed": len(errors),
            "scan_path": str(dir_path),
            "max_depth": max_depth,
            "recursive": recursive,
            "ocr_engine": ocr_engine,
            "language": language
        }

        return {
            "success": True,
            "scanned_files": len(found_images),
            "ocr_results": ocr_results,
            "errors": errors,
            "summary": summary
        }

    except Exception as e:
        logger.error(f"Error in scan_directory_and_ocr: {e}")
        return {
            "success": False,
            "error": str(e),
            "scanned_files": 0,
            "ocr_results": {},
            "errors": []
        }

# Helper functions for waiting on task completion
def _wait_for_task_completion(task_id: str, timeout: int) -> Dict[str, Any]:
    """Wait for single task completion"""
    import time
    start_time = time.time()

    while time.time() - start_time < timeout:
        status = ocr_queue.get_task_status(task_id)
        if status and status["status"] in ["completed", "failed"]:
            return {
                "success": True,
                "mode": "single_file",
                "task_id": task_id,
                "result": status
            }
        time.sleep(1)

    return {
        "success": False,
        "error": "Timeout waiting for task completion",
        "task_id": task_id
    }

def _wait_for_group_completion(group_id: str, timeout: int) -> Dict[str, Any]:
    """Wait for batch group completion"""
    import time
    start_time = time.time()

    while time.time() - start_time < timeout:
        status = ocr_queue.get_batch_status(group_id)
        if status:
            total = status["total_tasks"]
            completed = status["completed"] + status["failed"]

            if completed >= total:
                return {
                    "success": True,
                    "mode": "batch",
                    "group_id": group_id,
                    "result": status
                }
        time.sleep(2)

    return {
        "success": False,
        "error": "Timeout waiting for batch completion",
        "group_id": group_id
    }

def _wait_for_batch_completion(batch_info: Dict[str, Any], timeout: int) -> Dict[str, Any]:
    """Wait for 2D batch completion"""
    import time
    start_time = time.time()

    group_results = []

    while time.time() - start_time < timeout:
        all_completed = True

        for group in batch_info["groups"]:
            group_id = group["group_id"]
            status = ocr_queue.get_batch_status(group_id)

            if status:
                total = status["total_tasks"]
                completed = status["completed"] + status["failed"]

                if completed < total:
                    all_completed = False
                    break
                else:
                    # Update group result
                    group["status"] = status

        if all_completed:
            return {
                "success": True,
                "mode": "2d_queue",
                "batch_id": batch_info["batch_id"],
                "groups": batch_info["groups"],
                "total_files": batch_info["total_files"]
            }

        time.sleep(3)

    return {
        "success": False,
        "error": "Timeout waiting for 2D batch completion",
        "batch_id": batch_info["batch_id"],
        "partial_results": batch_info["groups"]
    }

@mcp.tool()
def analyze_image_pixels(
    image_path: str,
    region: str = "full",
    regions: Optional[List[str]] = None,
    deduplicate: bool = True,
    sample_size: int = 1000,
    sampling_strategy: str = "random",
    output_format: str = "hex",
    include_coordinates: bool = False,
    include_frequency: bool = True,
    color_tolerance: int = 0,
    max_dimension: Optional[int] = None,
    export_to_file: bool = True,
    export_path: Optional[str] = None,
    export_format: str = "json",
    export_filename: Optional[str] = None,
    return_inline_data: bool = False
) -> Dict[str, Any]:
    """
    Advanced image pixel analysis with customizable export options for AI workflows.

    IMPORTANT FOR AI: This function exports pixel data to files by default.
    DO NOT read pixel arrays token by token! Instead, use the Read tool to read the exported file.
    The response provides a file path - use it to efficiently access all pixel data at once.

    Export Format Guide:
    - "json": Full structured data with metadata (read with json.load())
    - "json-compact": Smaller JSON file without indentation
    - "csv": Simple spreadsheet format - columns: hex, r, g, b (or x, y, hex, r, g, b)
    - "txt": Human-readable plain text format
    - "python": Python dict literal (can be imported directly with eval())

    Example workflow for AI:
    1. Call analyze_image_pixels() with desired parameters
    2. Receive response with export_info.file_path
    3. Use Read tool to read the file at export_info.file_path
    4. Process pixel data efficiently in one operation

    Args:
        image_path: Path to the image file

        region: Single region to extract (ignored if regions is set):
            - "full": Entire image (default)
            - "left", "center", "right": Third columns
            - "top-left", "top-center", "top-right": 3x3 grid cells
            - "middle-left", "middle-center", "middle-right": 3x3 grid cells
            - "bottom-left", "bottom-center", "bottom-right": 3x3 grid cells
            - "custom:x1,y1,x2,y2": Custom rectangle (e.g., "custom:100,200,400,600")

        regions: List of regions for batch extraction (e.g., ["left", "center", "right"])

        deduplicate: Remove duplicate colors (default: True)

        sample_size: Number of pixel samples to extract (default: 1000)

        sampling_strategy: How to sample pixels:
            - "random": Random sampling (default)
            - "uniform": Uniformly distributed sampling
            - "grid": Grid-based sampling
            - "edge": Edge-priority sampling

        output_format: Output format - "hex", "rgb", "both"

        include_coordinates: Include pixel coordinates in output (default: False)

        include_frequency: Include color frequency statistics (default: True)

        color_tolerance: Merge similar colors within tolerance (0-255, default: 0 = exact match)

        max_dimension: Resize image if larger than this (for performance, default: None)

        export_to_file: Export results to file (default: True)

        export_path: Custom directory path for export (default: {project_root}/.cache/file_processor)
            - Absolute path: "D:/my_exports" or "/home/user/exports"
            - Relative to project: "exports/pixels" -> {project_root}/exports/pixels
            - Use "default" or None for default cache directory

        export_format: Export file format (default: "json"):
            - "json": JSON format with indentation
            - "json-compact": Compact JSON (no indentation, smaller file)
            - "csv": CSV format (colors only, good for spreadsheets)
            - "txt": Plain text format (human-readable)
            - "python": Python dict literal (can be imported directly)

        export_filename: Custom filename (default: auto-generated with timestamp)
            - Without extension: "my_colors" -> "my_colors.json"
            - With extension: "colors.csv" (extension must match export_format)

        return_inline_data: Return full data inline even if exported (default: False)
            - False: Return summary with file path (AI-friendly, avoids large responses)
            - True: Return all data inline + file path (may cause timeout for large data)

    Returns:
        AI-optimized response with file path and summary or full data based on return_inline_data
    """
    try:
        logger.info(f"Analyzing image pixels: {image_path}, region={region}, strategy={sampling_strategy}")

        try:
            from PIL import Image
            import numpy as np
            from collections import Counter
            IMAGE_PROCESSING_AVAILABLE = True
        except ImportError:
            return {
                "error": "Image processing not available. Please install: pip install Pillow numpy"
            }

        # Validate parameters
        if sample_size <= 0:
            return {"error": "sample_size must be greater than 0"}

        if color_tolerance < 0 or color_tolerance > 255:
            return {"error": "color_tolerance must be between 0 and 255"}

        if sampling_strategy not in ["random", "uniform", "grid", "edge"]:
            return {"error": f"Invalid sampling_strategy: {sampling_strategy}. Valid options: random, uniform, grid, edge"}

        # Validate export format
        valid_export_formats = ["json", "json-compact", "csv", "txt", "python"]
        if export_format not in valid_export_formats:
            return {"error": f"Invalid export_format: {export_format}. Valid options: {', '.join(valid_export_formats)}"}

        # Resolve image path
        image_file = Path(image_path)
        if not image_file.is_absolute():
            image_file = Path.cwd() / image_file

        if not image_file.exists():
            return {"error": f"Image file not found: {image_file}"}

        # Setup export directory
        try:
            from constants import FileProcessorConstants
            project_root = FileProcessorConstants.get_project_root()
        except:
            project_root = Path(__file__).parent.parent.parent.parent

        # Determine export directory
        if export_path:
            if export_path == "default":
                export_dir = project_root / ".cache" / "file_processor"
            else:
                export_dir = Path(export_path)
                if not export_dir.is_absolute():
                    # Relative to project root
                    export_dir = project_root / export_dir
        else:
            export_dir = project_root / ".cache" / "file_processor"

        export_dir.mkdir(parents=True, exist_ok=True)

        try:
            img = Image.open(image_file).convert("RGB")

            # Resize if needed for performance
            original_size = img.size
            if max_dimension and (img.width > max_dimension or img.height > max_dimension):
                ratio = min(max_dimension / img.width, max_dimension / img.height)
                new_size = (int(img.width * ratio), int(img.height * ratio))
                img = img.resize(new_size, Image.Resampling.LANCZOS)
                logger.info(f"Resized image from {original_size} to {img.size}")

            img_array = np.array(img)
            height, width, channels = img_array.shape

            # Helper function to extract pixels from a region
            def extract_region_pixels(region_spec: str, img_arr: np.ndarray) -> Dict[str, Any]:
                h, w = img_arr.shape[:2]
                region_name = region_spec.lower()

                # Parse custom region
                if region_name.startswith("custom:"):
                    try:
                        coords = region_name.split(":", 1)[1]
                        x1, y1, x2, y2 = map(int, coords.split(","))
                        # Clamp to image bounds
                        x1, y1 = max(0, x1), max(0, y1)
                        x2, y2 = min(w, x2), min(h, y2)
                        if x1 >= x2 or y1 >= y2:
                            return {"error": f"Invalid custom region coordinates: {coords}"}
                        region_coords = (x1, y1, x2, y2)
                    except:
                        return {"error": f"Invalid custom region format. Use: custom:x1,y1,x2,y2"}

                # Three-column regions
                elif region_name == "left":
                    region_coords = (0, 0, w // 3, h)
                elif region_name == "center":
                    region_coords = (w // 3, 0, 2 * w // 3, h)
                elif region_name == "right":
                    region_coords = (2 * w // 3, 0, w, h)

                # Nine-grid regions
                elif region_name == "top-left":
                    region_coords = (0, 0, w // 3, h // 3)
                elif region_name == "top-center":
                    region_coords = (w // 3, 0, 2 * w // 3, h // 3)
                elif region_name == "top-right":
                    region_coords = (2 * w // 3, 0, w, h // 3)
                elif region_name == "middle-left":
                    region_coords = (0, h // 3, w // 3, 2 * h // 3)
                elif region_name == "middle-center":
                    region_coords = (w // 3, h // 3, 2 * w // 3, 2 * h // 3)
                elif region_name == "middle-right":
                    region_coords = (2 * w // 3, h // 3, w, 2 * h // 3)
                elif region_name == "bottom-left":
                    region_coords = (0, 2 * h // 3, w // 3, h)
                elif region_name == "bottom-center":
                    region_coords = (w // 3, 2 * h // 3, 2 * w // 3, h)
                elif region_name == "bottom-right":
                    region_coords = (2 * w // 3, 2 * h // 3, w, h)

                # Full image
                elif region_name == "full":
                    region_coords = (0, 0, w, h)
                else:
                    return {"error": f"Invalid region: {region_spec}"}

                # Extract region
                x1, y1, x2, y2 = region_coords
                region_array = img_arr[y1:y2, x1:x2]
                region_h, region_w = region_array.shape[:2]

                return {
                    "coords": region_coords,
                    "array": region_array,
                    "name": region_name,
                    "width": region_w,
                    "height": region_h
                }

            # Helper function to apply sampling strategy
            def apply_sampling(pixels: np.ndarray, coords_list: List, strategy: str, n: int) -> Tuple[np.ndarray, List]:
                available = len(pixels)
                if available <= n:
                    return pixels, coords_list

                if strategy == "random":
                    indices = np.random.choice(available, n, replace=False)
                    return pixels[indices], [coords_list[i] for i in indices] if coords_list else []

                elif strategy == "uniform":
                    step = available / n
                    indices = [int(i * step) for i in range(n)]
                    return pixels[indices], [coords_list[i] for i in indices] if coords_list else []

                elif strategy == "grid":
                    # Sample in grid pattern
                    indices = np.linspace(0, available - 1, n, dtype=int)
                    return pixels[indices], [coords_list[i] for i in indices] if coords_list else []

                elif strategy == "edge":
                    # Prioritize edge pixels (first and last pixels in sorted order)
                    half = n // 2
                    indices = list(range(half)) + list(range(available - half, available))
                    indices = indices[:n]
                    return pixels[indices], [coords_list[i] for i in indices] if coords_list else []

                return pixels[:n], coords_list[:n] if coords_list else []

            # Helper function to merge similar colors
            def merge_similar_colors(pixels: np.ndarray, tolerance: int) -> Tuple[np.ndarray, Dict]:
                if tolerance == 0:
                    return pixels, {}

                # Round colors to tolerance buckets
                bucketed = (pixels // (tolerance + 1)) * (tolerance + 1)
                unique_buckets, inverse = np.unique(bucketed, axis=0, return_inverse=True)

                # Map original colors to buckets
                color_map = {}
                for i, bucket in enumerate(unique_buckets):
                    bucket_tuple = tuple(bucket)
                    original_colors = pixels[inverse == i]
                    avg_color = original_colors.mean(axis=0).astype(int)
                    color_map[bucket_tuple] = tuple(avg_color)

                return unique_buckets, color_map

            # Process regions
            region_list = regions if regions else [region]
            results = []

            for reg_spec in region_list:
                reg_data = extract_region_pixels(reg_spec, img_array)

                if "error" in reg_data:
                    results.append({"region": reg_spec, "error": reg_data["error"]})
                    continue

                region_array = reg_data["array"]

                # Flatten pixels and get coordinates if needed
                if include_coordinates:
                    coords = []
                    pixels_list = []
                    for y in range(reg_data["height"]):
                        for x in range(reg_data["width"]):
                            pixels_list.append(region_array[y, x])
                            coords.append((x + reg_data["coords"][0], y + reg_data["coords"][1]))
                    pixels = np.array(pixels_list)
                else:
                    pixels = region_array.reshape(-1, 3)
                    coords = []

                total_pixels = len(pixels)

                # Color frequency before deduplication
                if include_frequency:
                    color_freq = Counter([tuple(p) for p in pixels])

                # Apply color tolerance merging
                if color_tolerance > 0:
                    merged_pixels, color_map = merge_similar_colors(pixels, color_tolerance)
                    if include_frequency and color_map:
                        # Update frequency with merged colors
                        new_freq = Counter()
                        for orig_color, count in color_freq.items():
                            bucket = tuple((np.array(orig_color) // (color_tolerance + 1)) * (color_tolerance + 1))
                            merged_color = color_map.get(bucket, orig_color)
                            new_freq[merged_color] += count
                        color_freq = new_freq
                    pixels = merged_pixels

                # Deduplication
                if deduplicate:
                    unique_pixels, unique_indices = np.unique(pixels, axis=0, return_index=True)
                    pixels_to_sample = unique_pixels
                    coords_to_sample = [coords[i] for i in unique_indices] if coords else []
                    dedupe_count = len(unique_pixels)
                else:
                    pixels_to_sample = pixels
                    coords_to_sample = coords
                    dedupe_count = None

                # Apply sampling strategy
                sampled_pixels, sampled_coords = apply_sampling(
                    pixels_to_sample, coords_to_sample, sampling_strategy, sample_size
                )
                actual_sample_size = len(sampled_pixels)

                # Build region result
                region_result = {
                    "region": reg_data["name"],
                    "region_info": {
                        "coordinates": {
                            "x1": int(reg_data["coords"][0]),
                            "y1": int(reg_data["coords"][1]),
                            "x2": int(reg_data["coords"][2]),
                            "y2": int(reg_data["coords"][3])
                        },
                        "width": int(reg_data["width"]),
                        "height": int(reg_data["height"]),
                        "total_pixels": int(total_pixels)
                    },
                    "processing_info": {
                        "deduplicated": deduplicate,
                        "unique_colors": int(dedupe_count) if dedupe_count is not None else None,
                        "color_tolerance": color_tolerance,
                        "sampling_strategy": sampling_strategy,
                        "requested_sample_size": sample_size,
                        "actual_sample_size": int(actual_sample_size)
                    }
                }

                # Format output
                if output_format in ["hex", "both"]:
                    if include_coordinates:
                        region_result["hex_pixels"] = [
                            {"color": f"#{r:02x}{g:02x}{b:02x}", "x": x, "y": y}
                            for (r, g, b), (x, y) in zip(sampled_pixels, sampled_coords)
                        ]
                    else:
                        region_result["hex_colors"] = [f"#{r:02x}{g:02x}{b:02x}" for r, g, b in sampled_pixels]

                if output_format in ["rgb", "both"]:
                    if include_coordinates:
                        region_result["rgb_pixels"] = [
                            {"color": [int(r), int(g), int(b)], "x": x, "y": y}
                            for (r, g, b), (x, y) in zip(sampled_pixels, sampled_coords)
                        ]
                    else:
                        region_result["rgb_colors"] = [[int(r), int(g), int(b)] for r, g, b in sampled_pixels]

                # Add frequency statistics
                if include_frequency:
                    sorted_freq = sorted(color_freq.items(), key=lambda x: x[1], reverse=True)
                    region_result["color_frequency"] = {
                        "unique_colors": len(color_freq),
                        "most_frequent": [
                            {
                                "rgb": [int(r), int(g), int(b)],
                                "hex": f"#{r:02x}{g:02x}{b:02x}",
                                "count": int(count),
                                "percentage": round(count / total_pixels * 100, 2)
                            }
                            for (r, g, b), count in sorted_freq[:10]
                        ]
                    }

                results.append(region_result)

            # Build final result
            final_result = {
                "success": True,
                "file_path": str(image_file),
                "image_info": {
                    "original_size": {"width": original_size[0], "height": original_size[1]} if max_dimension else None,
                    "processed_size": {"width": width, "height": height},
                    "channels": channels,
                    "format": img.format,
                    "mode": img.mode
                },
                "regions": results if len(results) > 1 else results[0]
            }

            # Export to file if requested
            if export_to_file:
                # Determine filename
                if export_filename:
                    # Check if filename has extension
                    if '.' in export_filename:
                        filename = export_filename
                    else:
                        # Add extension based on export_format
                        ext_map = {
                            "json": ".json",
                            "json-compact": ".json",
                            "csv": ".csv",
                            "txt": ".txt",
                            "python": ".py"
                        }
                        filename = export_filename + ext_map.get(export_format, ".txt")
                else:
                    timestamp = time.strftime("%Y%m%d_%H%M%S")
                    ext_map = {
                        "json": ".json",
                        "json-compact": ".json",
                        "csv": ".csv",
                        "txt": ".txt",
                        "python": ".py"
                    }
                    filename = f"pixels_{timestamp}{ext_map.get(export_format, '.txt')}"

                file_path = export_dir / filename

                # Export based on format
                try:
                    if export_format == "json":
                        with open(file_path, 'w', encoding='utf-8') as f:
                            json.dump(final_result, f, indent=2, ensure_ascii=False)

                    elif export_format == "json-compact":
                        with open(file_path, 'w', encoding='utf-8') as f:
                            json.dump(final_result, f, ensure_ascii=False)

                    elif export_format == "csv":
                        import csv
                        with open(file_path, 'w', newline='', encoding='utf-8') as f:
                            writer = csv.writer(f)
                            # Write header
                            if include_coordinates:
                                writer.writerow(['x', 'y', 'hex', 'r', 'g', 'b'])
                            else:
                                writer.writerow(['hex', 'r', 'g', 'b'])

                            # Write data from regions
                            region_data = final_result["regions"]
                            if not isinstance(region_data, list):
                                region_data = [region_data]

                            for reg in region_data:
                                if "hex_colors" in reg:
                                    for hex_color in reg["hex_colors"]:
                                        r = int(hex_color[1:3], 16)
                                        g = int(hex_color[3:5], 16)
                                        b = int(hex_color[5:7], 16)
                                        writer.writerow([hex_color, r, g, b])
                                elif "hex_pixels" in reg:
                                    for pixel in reg["hex_pixels"]:
                                        hex_color = pixel["color"]
                                        r = int(hex_color[1:3], 16)
                                        g = int(hex_color[3:5], 16)
                                        b = int(hex_color[5:7], 16)
                                        writer.writerow([pixel["x"], pixel["y"], hex_color, r, g, b])

                    elif export_format == "txt":
                        with open(file_path, 'w', encoding='utf-8') as f:
                            f.write("=" * 60 + "\n")
                            f.write("Image Pixel Analysis Results\n")
                            f.write("=" * 60 + "\n\n")
                            f.write(f"Image: {image_file}\n")
                            f.write(f"Processed Size: {width}x{height}\n\n")

                            region_data = final_result["regions"]
                            if not isinstance(region_data, list):
                                region_data = [region_data]

                            for reg in region_data:
                                f.write(f"\n--- Region: {reg['region']} ---\n")
                                f.write(f"Samples: {reg['processing_info']['actual_sample_size']}\n\n")

                                if "hex_colors" in reg:
                                    for i, hex_color in enumerate(reg["hex_colors"], 1):
                                        f.write(f"{i}. {hex_color}\n")
                                elif "hex_pixels" in reg:
                                    for i, pixel in enumerate(reg["hex_pixels"], 1):
                                        f.write(f"{i}. ({pixel['x']}, {pixel['y']}) = {pixel['color']}\n")

                    elif export_format == "python":
                        with open(file_path, 'w', encoding='utf-8') as f:
                            f.write("# Auto-generated pixel data\n")
                            f.write(f"# Generated: {time.strftime('%Y-%m-%d %H:%M:%S')}\n\n")
                            f.write("pixel_data = ")
                            f.write(repr(final_result))

                    file_size = file_path.stat().st_size

                    # Create AI-friendly response
                    if return_inline_data:
                        # Return full data + file info
                        final_result["export_info"] = {
                            "exported": True,
                            "file_path": str(file_path),
                            "file_name": filename,
                            "file_size_bytes": file_size,
                            "file_size_kb": round(file_size / 1024, 2),
                            "export_directory": str(export_dir),
                            "export_format": export_format,
                            "data_included_inline": True
                        }
                        logger.info(f"Exported to: {file_path} (data returned inline)")
                        return final_result
                    else:
                        # Return summary only (AI-friendly, fast response)
                        region_data = final_result["regions"]
                        if not isinstance(region_data, list):
                            region_data = [region_data]

                        summary_response = {
                            "success": True,
                            "message": "Analysis complete. Data exported to file. AI should read the file for pixel data.",
                            "export_info": {
                                "exported": True,
                                "file_path": str(file_path),
                                "file_name": filename,
                                "file_size_bytes": file_size,
                                "file_size_kb": round(file_size / 1024, 2),
                                "export_directory": str(export_dir),
                                "export_format": export_format,
                                "data_included_inline": False
                            },
                            "image_info": final_result["image_info"],
                            "summary": {
                                "regions_processed": len(region_data),
                                "regions": [
                                    {
                                        "region": reg["region"],
                                        "total_pixels": reg["region_info"]["total_pixels"],
                                        "samples_extracted": reg["processing_info"]["actual_sample_size"],
                                        "unique_colors": reg["processing_info"].get("unique_colors"),
                                        "top_color": reg.get("color_frequency", {}).get("most_frequent", [{}])[0] if "color_frequency" in reg else None
                                    }
                                    for reg in region_data
                                ]
                            },
                            "usage_tip": f"To access pixel data, read the file: {file_path}"
                        }

                        logger.info(f"Exported to: {file_path} (summary only returned)")
                        return summary_response

                except Exception as export_error:
                    logger.error(f"Export failed: {export_error}")
                    final_result["export_info"] = {
                        "exported": False,
                        "error": str(export_error)
                    }
                    return final_result

            # No export, return inline data
            logger.info(f"Image pixel analysis completed: {len(region_list)} region(s) processed (no export)")
            return final_result

        except Exception as img_error:
            return {"error": f"Failed to process image: {str(img_error)}"}

    except Exception as e:
        error_msg = f"Image pixel analysis failed: {str(e)}"
        logger.error(error_msg)
        return {"error": error_msg}

# ==================== IMAGE PROCESSING TOOLS ====================

@mcp.tool()
def split_image_equal(
    image_path: str,
    count: int,
    direction: str = "vertical",
    output_dir: Optional[str] = None,
    name_pattern: str = "part_{index}"
) -> Dict[str, Any]:
    """
    Split image into equal parts (auto-calculate part size)

    Args:
        image_path: Input image path
        count: Number of equal parts
        direction: Split direction ('horizontal' or 'vertical', default: 'vertical')
        output_dir: Output directory (default: same as input)
        name_pattern: Naming pattern with {index} placeholder

    Returns:
        Result dictionary with output files list

    Example:
        Split a 56x560 image into 10 equal parts vertically:
        split_image_equal("sprite.png", 10, "vertical")
        Result: 10 images of 56x56 each
    """
    try:
        from image_tools import image_tools
        return image_tools.split_image_equal(
            image_path, count, direction, output_dir, name_pattern
        )
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def split_image_custom(
    image_path: str,
    split_points: List[int],
    direction: str = "vertical",
    output_dir: Optional[str] = None,
    name_pattern: str = "part_{index}"
) -> Dict[str, Any]:
    """
    Split image at custom pixel positions

    Args:
        image_path: Input image path
        split_points: List of split positions in pixels
        direction: Split direction ('horizontal' or 'vertical', default: 'vertical')
        output_dir: Output directory (default: same as input)
        name_pattern: Naming pattern with {index} placeholder

    Returns:
        Result dictionary with output files list

    Example:
        split_points=[100, 250, 400] creates parts at:
        Part 0: 0-100px, Part 1: 100-250px, Part 2: 250-400px, Part 3: 400-end
    """
    try:
        from image_tools import image_tools
        return image_tools.split_image_custom(
            image_path, split_points, direction, output_dir, name_pattern
        )
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def create_image_grid(
    image_paths: List[str],
    cols: int,
    output_path: str,
    spacing: int = 0,
    background_color: str = "white",
    cell_width: Optional[int] = None,
    cell_height: Optional[int] = None,
    resize_mode: str = "fit"
) -> Dict[str, Any]:
    """
    Create image grid/collage from multiple images

    Args:
        image_paths: List of image paths to arrange
        cols: Number of columns
        output_path: Output path for the grid image
        spacing: Space between images in pixels (default: 0)
        background_color: Background color (default: "white")
        cell_width: Fixed cell width (optional, auto-detect from images if not set)
        cell_height: Fixed cell height (optional, auto-detect from images if not set)
        resize_mode: How to resize images in cells:
            - 'fit': Fit inside cell, maintain aspect (default)
            - 'fill': Fill cell, crop if needed, maintain aspect
            - 'stretch': Stretch to fill cell exactly

    Returns:
        Result dictionary with grid info

    Example:
        Create 3-column grid from 9 images:
        create_image_grid(["img1.png", "img2.png", ...], 3, "grid.png")
    """
    try:
        from image_tools import image_tools
        return image_tools.create_image_grid(
            image_paths, cols, output_path, spacing,
            background_color, cell_width, cell_height, resize_mode
        )
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def crop_image(
    image_path: str,
    x: int,
    y: int,
    width: int,
    height: int,
    output_path: Optional[str] = None
) -> Dict[str, Any]:
    """
    Crop image to specified rectangle

    Args:
        image_path: Input image path
        x: Left coordinate
        y: Top coordinate
        width: Crop width
        height: Crop height
        output_path: Output path (optional)

    Returns:
        Result dictionary with output path and metadata
    """
    try:
        from image_tools import image_tools
        return image_tools.crop_image(image_path, x, y, width, height, output_path)
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def split_sprite_sheet(
    image_path: str,
    sprite_width: int,
    sprite_height: int,
    output_dir: Optional[str] = None,
    name_pattern: str = "sprite_{index}",
    direction: str = "vertical"
) -> Dict[str, Any]:
    """
    Split sprite sheet into individual sprites

    Args:
        image_path: Input sprite sheet path
        sprite_width: Width of each sprite
        sprite_height: Height of each sprite
        output_dir: Output directory (default: same as input)
        name_pattern: Naming pattern with {index} placeholder
        direction: Split direction ('horizontal' or 'vertical', default: 'vertical')

    Returns:
        Result dictionary with sprite files list

    Example:
        Split a 80x400 vertical sprite sheet (5 sprites of 80x80 each):
        split_sprite_sheet("sprite.png", 80, 80, direction="vertical")
    """
    try:
        from image_tools import image_tools
        return image_tools.split_sprite_sheet(
            image_path, sprite_width, sprite_height,
            output_dir, name_pattern, direction
        )
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def split_image_grid(
    image_path: str,
    rows: int,
    cols: int,
    output_dir: Optional[str] = None,
    name_pattern: str = "tile_{row}_{col}"
) -> Dict[str, Any]:
    """
    Split image into grid (rows x cols)

    Args:
        image_path: Input image path
        rows: Number of rows
        cols: Number of columns
        output_dir: Output directory (default: same as input)
        name_pattern: Naming pattern with {row} and {col} placeholders

    Returns:
        Result dictionary with output files list
    """
    try:
        from image_tools import image_tools
        return image_tools.split_image_grid(image_path, rows, cols, output_dir, name_pattern)
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def resize_image(
    image_path: str,
    width: Optional[int] = None,
    height: Optional[int] = None,
    max_size: Optional[int] = None,
    keep_aspect: bool = True,
    resample: str = "LANCZOS",
    output_path: Optional[str] = None
) -> Dict[str, Any]:
    """
    Resize image with various options

    Args:
        image_path: Input image path
        width: Target width
        height: Target height
        max_size: Max dimension (used if width/height not specified)
        keep_aspect: Keep aspect ratio (default: True)
        resample: Resampling method (LANCZOS, BILINEAR, BICUBIC, NEAREST)
        output_path: Output path (optional)

    Returns:
        Result dictionary
    """
    try:
        from image_tools import image_tools
        return image_tools.resize_image(
            image_path, width, height, max_size,
            keep_aspect, resample, output_path
        )
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def rotate_image(
    image_path: str,
    angle: float,
    expand: bool = True,
    fill_color: str = "white",
    output_path: Optional[str] = None
) -> Dict[str, Any]:
    """
    Rotate image by angle

    Args:
        image_path: Input image path
        angle: Rotation angle in degrees (counter-clockwise)
        expand: Expand canvas to fit rotated image (default: True)
        fill_color: Background color for expanded area
        output_path: Output path (optional)

    Returns:
        Result dictionary
    """
    try:
        from image_tools import image_tools
        return image_tools.rotate_image(image_path, angle, expand, fill_color, output_path)
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def flip_image(
    image_path: str,
    direction: str = "horizontal",
    output_path: Optional[str] = None
) -> Dict[str, Any]:
    """
    Flip image horizontally or vertically

    Args:
        image_path: Input image path
        direction: 'horizontal' or 'vertical'
        output_path: Output path (optional)

    Returns:
        Result dictionary
    """
    try:
        from image_tools import image_tools
        return image_tools.flip_image(image_path, direction, output_path)
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def compress_image(
    image_path: str,
    quality: int = 85,
    max_width: Optional[int] = None,
    max_height: Optional[int] = None,
    output_format: Optional[str] = None,
    output_path: Optional[str] = None
) -> Dict[str, Any]:
    """
    Compress image with quality and size control

    Args:
        image_path: Input image path
        quality: JPEG quality (1-100, default: 85)
        max_width: Maximum width (optional)
        max_height: Maximum height (optional)
        output_format: Output format (JPEG, PNG, WEBP, optional)
        output_path: Output path (optional)

    Returns:
        Result dictionary with compression ratio
    """
    try:
        from image_tools import image_tools
        return image_tools.compress_image(
            image_path, quality, max_width, max_height,
            output_format, output_path
        )
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def merge_images_horizontal(
    image_paths: List[str],
    output_path: str,
    spacing: int = 0,
    align: str = "center"
) -> Dict[str, Any]:
    """
    Merge images horizontally

    Args:
        image_paths: List of image paths to merge
        output_path: Output path for merged image
        spacing: Space between images in pixels (default: 0)
        align: Vertical alignment ('top', 'center', 'bottom', default: 'center')

    Returns:
        Result dictionary
    """
    try:
        from image_tools import image_tools
        return image_tools.merge_images_horizontal(image_paths, output_path, spacing, align)
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def merge_images_vertical(
    image_paths: List[str],
    output_path: str,
    spacing: int = 0,
    align: str = "center"
) -> Dict[str, Any]:
    """
    Merge images vertically

    Args:
        image_paths: List of image paths to merge
        output_path: Output path for merged image
        spacing: Space between images in pixels (default: 0)
        align: Horizontal alignment ('left', 'center', 'right', default: 'center')

    Returns:
        Result dictionary
    """
    try:
        from image_tools import image_tools
        return image_tools.merge_images_vertical(image_paths, output_path, spacing, align)
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def add_text_to_image(
    image_path: str,
    text: str,
    position: Tuple[int, int] = (10, 10),
    font_size: int = 24,
    font_color: str = "black",
    font_path: Optional[str] = None,
    background_color: Optional[str] = None,
    output_path: Optional[str] = None
) -> Dict[str, Any]:
    """
    Add text overlay to image

    Args:
        image_path: Input image path
        text: Text to add
        position: (x, y) position tuple (default: (10, 10))
        font_size: Font size in pixels (default: 24)
        font_color: Text color (default: "black")
        font_path: Custom font file path (TTF, optional)
        background_color: Text background color (optional)
        output_path: Output path (optional)

    Returns:
        Result dictionary
    """
    try:
        from image_tools import image_tools
        return image_tools.add_text_to_image(
            image_path, text, position, font_size,
            font_color, font_path, background_color, output_path
        )
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def apply_image_filter(
    image_path: str,
    filter_type: str,
    intensity: float = 1.0,
    output_path: Optional[str] = None
) -> Dict[str, Any]:
    """
    Apply image filter/effect

    Args:
        image_path: Input image path
        filter_type: Filter type (blur, sharpen, brightness, contrast, grayscale, sepia)
        intensity: Filter intensity (0.0 - 2.0, default: 1.0)
        output_path: Output path (optional)

    Returns:
        Result dictionary

    Example filters:
        - blur: Gaussian blur
        - sharpen: Sharpen image
        - brightness: Adjust brightness
        - contrast: Adjust contrast
        - grayscale: Convert to grayscale
        - sepia: Sepia tone effect
    """
    try:
        from image_tools import image_tools
        return image_tools.apply_filter(image_path, filter_type, intensity, output_path)
    except Exception as e:
        return {"error": str(e)}

# ==================== CODE SCANNING TOOLS ====================

@mcp.tool()
def scan_code_for_chinese(
    directory: str,
    extensions: Optional[List[str]] = None,
    max_depth: int = -1,
    output_format: str = "text"
) -> Union[str, Dict[str, Any]]:
    """
    Scan code files for Chinese characters in a directory.

    This tool recursively scans code files and detects Chinese characters,
    useful for code cleanup, internationalization audits, or finding hardcoded strings.

    Args:
        directory: Root directory to scan
        extensions: List of file extensions to scan (e.g., ['py', 'js', 'dart'])
                   If None, scans all default code file types:
                   - JavaScript/TypeScript: .js, .jsx, .ts, .tsx, .mjs, .cjs
                   - Python: .py, .pyw, .pyi
                   - Dart: .dart
                   - Vue: .vue
                   - Java/Kotlin: .java, .kt, .kts
                   - Go: .go
                   - Shell: .sh, .bash, .zsh, .fish, .ps1, .psm1, .psd1
                   - C/C++: .c, .cpp, .cc, .h, .hpp
                   - C#: .cs
                   - And many more...
        max_depth: Maximum directory depth to scan (-1 for unlimited, default: -1)
        output_format: Output format - 'text' for readable format or 'json' for structured data

    Returns:
        Scan results with statistics and locations of Chinese characters.
        Format depends on output_format parameter.

    Excluded directories:
        node_modules, .git, dist, build, __pycache__, venv, target, vendor, etc.

    Example usage:
        # Scan all code files
        scan_code_for_chinese("/path/to/project")

        # Scan only Python files
        scan_code_for_chinese("/path/to/project", extensions=["py"])

        # Scan with depth limit
        scan_code_for_chinese("/path/to/project", max_depth=3)

        # Get JSON output
        scan_code_for_chinese("/path/to/project", output_format="json")
    """
    try:
        if not CODE_SCANNER_AVAILABLE:
            error_msg = "Code scanner is not available. Please check installation."
            logger.error(error_msg)
            if output_format == "json":
                return {"success": False, "error": error_msg}
            return f"Error: {error_msg}"

        # Scan directory
        results = code_scanner.scan_directory(
            directory=directory,
            extensions=extensions,
            max_depth=max_depth
        )

        # Format output
        if output_format == "json":
            return results
        else:
            return code_scanner.format_results_text(results)

    except Exception as e:
        error_msg = f"Code scanning failed: {str(e)}"
        logger.error(error_msg)
        if output_format == "json":
            return {"success": False, "error": error_msg}
        return f"Error: {error_msg}"

@mcp.tool()
def list_code_extensions() -> Dict[str, Any]:
    """
    List all supported code file extensions for Chinese character scanning.

    Returns:
        Dictionary containing all supported extensions organized by language/category.

    Example usage:
        list_code_extensions()
    """
    try:
        if not CODE_SCANNER_AVAILABLE:
            return {"error": "Code scanner is not available"}

        extensions_by_category = {
            "JavaScript/TypeScript": [".js", ".jsx", ".ts", ".tsx", ".mjs", ".cjs"],
            "Python": [".py", ".pyw", ".pyi"],
            "Dart": [".dart"],
            "Vue": [".vue"],
            "Java/Kotlin": [".java", ".kt", ".kts"],
            "Go": [".go"],
            "Shell Scripts": [".sh", ".bash", ".zsh", ".fish"],
            "PowerShell": [".ps1", ".psm1", ".psd1"],
            "C/C++": [".c", ".cpp", ".cc", ".cxx", ".h", ".hpp", ".hxx"],
            "C#": [".cs"],
            "Ruby": [".rb"],
            "PHP": [".php"],
            "Rust": [".rs"],
            "Swift": [".swift"],
            "Scala": [".scala"],
            "R": [".r", ".R"],
            "Perl": [".pl", ".pm"],
            "Lua": [".lua"],
            "SQL": [".sql"],
            "Web": [".html", ".htm", ".css", ".scss", ".sass", ".less"],
            "Config": [".json", ".yaml", ".yml", ".toml", ".ini", ".conf", ".config"],
            "Markdown": [".md", ".markdown"],
            "XML": [".xml"]
        }

        return {
            "success": True,
            "total_extensions": sum(len(exts) for exts in extensions_by_category.values()),
            "categories": extensions_by_category,
            "excluded_directories": list(code_scanner.EXCLUDE_DIRS),
            "excluded_files": list(code_scanner.EXCLUDE_FILES),
            "usage": "Use extensions parameter in scan_code_for_chinese() to filter by extension"
        }

    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def health_check() -> Dict[str, Any]:
    """Check server health and package availability"""
    try:
        status = {
            "server": "healthy",
            "packages": {},
            "capabilities": [],
            "converters": []
        }

        # Check package availability
        for format_type, packages in PackageManager.PACKAGE_MAPPING.items():
            format_status = {}
            for package in packages:
                try:
                    importlib.import_module(package.replace('-', '_').replace('python_', ''))
                    format_status[package] = "available"
                except ImportError:
                    format_status[package] = "missing"

            status["packages"][format_type] = format_status

            # Determine capabilities
            available_packages = sum(1 for s in format_status.values() if s == "available")
            if available_packages >= len(packages) // 2 + 1:
                status["capabilities"].append(format_type)

        # Check converter capabilities
        for conversion in converter.conversion_mappings:
            status["converters"].append(conversion)

        # Check OCR availability
        status["ocr"] = {
            "available": OCR_AVAILABLE,
            "engines": ocr_manager.get_available_engines() if OCR_AVAILABLE else [],
            "queue_system": ocr_queue is not None and ocr_queue.is_running if OCR_AVAILABLE else False,
            "smart_features": {
                "image_compression": True,
                "pdf_splitting": True,
                "queue_processing": True,
                "batch_support": True,
                "2d_queue_support": True
            } if OCR_AVAILABLE else {}
        }

        # Add OCR system stats if available
        if OCR_AVAILABLE and ocr_queue:
            try:
                ocr_stats = ocr_queue.get_system_stats()
                status["ocr"]["system_stats"] = ocr_stats
            except Exception as e:
                status["ocr"]["system_stats_error"] = str(e)

        # Add image pixel analysis capability
        status["image_pixel_analysis"] = {
            "available": True,
            "supported_regions": [
                "full", "left", "center", "right",
                "top-left", "top-center", "top-right",
                "middle-left", "middle-center", "middle-right",
                "bottom-left", "bottom-center", "bottom-right",
                "custom:x1,y1,x2,y2"
            ],
            "features": {
                "region_extraction": True,
                "batch_regions": True,
                "deduplication": True,
                "sampling_strategies": ["random", "uniform", "grid", "edge"],
                "color_tolerance": True,
                "coordinate_tracking": True,
                "frequency_statistics": True,
                "performance_optimization": True,
                "file_export": True
            },
            "output_formats": ["hex", "rgb", "both"],
            "exchange_directory": "{project_root}/.cache/file_processor"
        }

        # Add code scanner capability
        status["code_scanner"] = {
            "available": CODE_SCANNER_AVAILABLE,
            "features": {
                "chinese_detection": True,
                "recursive_scanning": True,
                "multi_encoding_support": True,
                "extension_filtering": True,
                "depth_limiting": True,
                "smart_exclusions": True
            },
            "supported_languages": [
                "JavaScript/TypeScript", "Python", "Dart", "Vue",
                "Java/Kotlin", "Go", "Shell", "PowerShell",
                "C/C++", "C#", "Ruby", "PHP", "Rust", "Swift",
                "Scala", "R", "Perl", "Lua", "SQL", "HTML/CSS"
            ],
            "total_extensions": 60,
            "output_formats": ["text", "json"]
        } if CODE_SCANNER_AVAILABLE else {"available": False}

        return status

    except Exception as e:
        return {"server": "error", "error": str(e)}

def _initialize_ocr_engines():
    """Initialize all available OCR engines on startup"""
    try:
        from ocr_engines import OCRManager

        # Check if this is a restart/reconnection scenario
        restart_file = Path("tmp_mcp_restart_marker")
        is_restart = restart_file.exists()

        if is_restart:
            logger.info("Detected MCP restart - handling reconnection scenario")
            restart_file.unlink()  # Remove restart marker
        else:
            logger.info("Fresh MCP startup - initializing OCR engines")

        ocr_manager = OCRManager()

        # Get available engines to check what's installed
        available_engines = ocr_manager.get_available_engines()
        logger.info(f"Available OCR engines: {available_engines}")

        # Try to initialize PaddleOCR first
        if 'paddle' in available_engines:
            try:
                if ocr_manager.initialize_paddle_ocr():
                    logger.info("PaddleOCR initialized successfully")
                else:
                    logger.warning("PaddleOCR initialization failed")
            except Exception as e:
                logger.warning(f"PaddleOCR initialization error: {e}")

        # Try to initialize CnOCR (with fallback models)
        if 'cnocr' in available_engines:
            model_types_to_try = ["general", "scene", "doc"]
            cnocr_initialized = False

            for model_type in model_types_to_try:
                try:
                    if ocr_manager.initialize_cnocr(model_type):
                        logger.info(f"CnOCR initialized successfully with {model_type} model")
                        cnocr_initialized = True
                        break
                except Exception as e:
                    logger.warning(f"CnOCR {model_type} model initialization error: {e}")

            if not cnocr_initialized:
                logger.warning("All CnOCR model initialization attempts failed")

        # Free OCR doesn't need initialization
        if 'free' in available_engines:
            logger.info("Free OCR service available")

        # Clean up any temporary files from previous sessions
        _cleanup_temp_files()

        logger.info("OCR engines initialization completed")

    except Exception as e:
        logger.error(f"OCR engines initialization failed: {e}")

def _normalize_file_path(file_path: str) -> str:
    """Normalize file path for better compatibility with Chinese characters and Windows paths"""
    try:
        import os
        # Convert to absolute path
        normalized = os.path.abspath(file_path)

        # Handle Windows path separators
        if os.name == 'nt':
            normalized = normalized.replace('/', '\\')

        # Ensure proper encoding for Chinese characters
        if isinstance(normalized, str):
            normalized = normalized.encode('utf-8').decode('utf-8')

        return normalized
    except Exception as e:
        logger.warning(f"Path normalization warning: {e}")
        return file_path

def _validate_image_file(file_path: str) -> bool:
    """Validate that image file exists and is accessible"""
    try:
        if not os.path.exists(file_path):
            return False

        if not os.path.isfile(file_path):
            return False

        # Check file size (should be > 0)
        if os.path.getsize(file_path) <= 0:
            return False

        # Check file extension
        valid_extensions = ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.gif', '.webp']
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext not in valid_extensions:
            logger.warning(f"Unusual image file extension: {file_ext}")

        return True
    except Exception as e:
        logger.error(f"File validation error: {e}")
        return False

def _cleanup_temp_files():
    """Clean up temporary files from previous sessions"""
    try:
        import tempfile
        import glob

        # Clean up temporary OCR files
        temp_dir = tempfile.gettempdir()
        temp_patterns = [
            "tmp_ocr_*",
            "tmp_mcp_*",
            "tmp_paddle_*",
            "tmp_cnocr_*"
        ]

        cleaned_count = 0
        for pattern in temp_patterns:
            temp_files = glob.glob(os.path.join(temp_dir, pattern))
            for temp_file in temp_files:
                try:
                    os.remove(temp_file)
                    cleaned_count += 1
                except OSError:
                    pass  # File might be in use or already deleted

        if cleaned_count > 0:
            logger.info(f"Cleaned up {cleaned_count} temporary files")

    except Exception as e:
        logger.warning(f"Temporary file cleanup warning: {e}")

def main():
    """Main MCP server function using FastMCP"""
    try:
        logger.info("Starting Document Parser MCP Server...")
        logger.info("Supported formats: PDF, XMind, Word, Excel, PowerPoint, Images, Text")

        # Initialize OCR engines proactively
        logger.info("Initializing OCR engines...")
        _initialize_ocr_engines()

        # Start FastMCP server
        logger.info("[FASTMCP] Starting FastMCP Document Parser server...")
        mcp.run()

    except KeyboardInterrupt:
        logger.info("[STOP] Server stopped by user")
        _create_restart_marker()
    except Exception as e:
        logger.error(f"[ERROR] Server startup failed: {e}")
        _create_restart_marker()
        raise

def _create_restart_marker():
    """Create restart marker for next startup"""
    try:
        restart_file = Path("tmp_mcp_restart_marker")
        restart_file.touch()
        logger.info("Created restart marker for next startup")
    except Exception as e:
        logger.warning(f"Failed to create restart marker: {e}")

if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        logger.error(f"Server error: {e}")
        sys.exit(1)